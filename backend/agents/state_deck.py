"""State Pitch Deck Generator — Oklahoma-style governor pitch.

Produces a 10-slide PPTX deck tailored for US state governors.
Format: Problem → Results → Live Pilot → How It Works → The Deal →
Economics → What Each Side Provides → Gated Roadmap → The Ask.

Adapted from the Oklahoma deck used for Gov Stitt.
"""

from __future__ import annotations

import logging
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from models.schemas import (
    CountryProfile, EducationAnalysis, Strategy,
    FinancialModel, FinancialAssumptions,
)
from services.llm import call_llm_plain
from config import OUTPUT_DIR
from config.rules_loader import get_esa_data

logger = logging.getLogger(__name__)

# Colour palette
DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0x6D, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xCC, 0x33, 0x33)

# ---------------------------------------------------------------------------
# Alpha reference data (confirmed by CFO)
# ---------------------------------------------------------------------------

ALPHA_RESULTS = {
    "avg_growth": "2.2x",
    "behind_growth": "4.9x",
    "love_school": "97%",
    "avg_sat": "1530",
    "ap_4_or_5": "94%",
    "lipscomb": "3x learning (Fall-Winter MAP)",
    "intervention_ela": "2x+",
    "intervention_math": "1.5-2x",
}

PRICING = {
    "intervention": 2_000,          # $/student/year for 1-hr block
    "mark_rober_science": 2_000,    # $/student/year for 1-hr block
    "timeback_pct": 20,             # % of per-pupil funding
}


# ---------------------------------------------------------------------------
# Prompt for state-specific research synthesis
# ---------------------------------------------------------------------------

STATE_DECK_PROMPT = """You are a senior director at Alpha Holdings preparing a governor pitch deck
for **{state}**. Using the data below, produce SPECIFIC NUMBERS for each slide.

State Data:
{state_data}

Education Data:
{education_data}

Strategy Summary:
{strategy_summary}

For each of the following, provide 3-5 bullet points with SPECIFIC DATA:

1. THE PROBLEM: {state}'s education rankings, NAEP scores, proficiency rates, graduation rate,
   per-pupil spending vs national average, total K-12 students, state education budget.
   Frame it as "this isn't a funding problem — it's a model problem."

2. WHAT {state_upper} GETS: What Alpha provides vs what the state provides (table format).
   Include alignment with governor's agenda and existing state policies.

3. ESA/VOUCHER OPPORTUNITY: If {state} has an ESA/voucher program, explain how Alpha's model
   works within it. ESA amount per student, total eligible students, how families can use
   ESA funds to access Alpha programs.

Return as structured markdown with ## headers for each section.
"""


async def generate_state_deck(
    state: str,
    country_profile: CountryProfile,
    education_analysis: EducationAnalysis,
    strategy: Strategy,
    financial_model: FinancialModel,
    assumptions: FinancialAssumptions,
) -> str:
    """Generate an Oklahoma-style state pitch deck (PPTX).

    Returns the path to the saved PPTX file.
    """
    logger.info("Generating state pitch deck for %s", state)

    # Get state-specific data
    esa_data = get_esa_data(state)
    per_pupil = country_profile.education.avg_public_spend_per_student or 11_000
    k12_students = country_profile.education.k12_enrolled or 500_000
    timeback_cost = round(per_pupil * PRICING["timeback_pct"] / 100)

    # Build state data context for LLM
    state_data_parts = []
    if country_profile.demographics.total_population:
        state_data_parts.append(f"Population: {country_profile.demographics.total_population:,.0f}")
    if k12_students:
        state_data_parts.append(f"K-12 students: {k12_students:,.0f}")
    if per_pupil:
        state_data_parts.append(f"Per-pupil spending: ${per_pupil:,.0f}")
    if country_profile.education.pisa_scores:
        state_data_parts.append(f"Test scores: {country_profile.education.pisa_scores}")
    if country_profile.political_context.head_of_state:
        state_data_parts.append(f"Governor: {country_profile.political_context.head_of_state}")
    if esa_data:
        state_data_parts.append(f"ESA program: {esa_data.get('program_name', 'Yes')}")
        state_data_parts.append(f"ESA amount: ${esa_data.get('esa_amount', 'N/A')}/student")

    education_parts = []
    if education_analysis.system_diagnosis.primary_pain_points:
        education_parts.append("Pain points: " + "; ".join(education_analysis.system_diagnosis.primary_pain_points[:5]))
    if education_analysis.system_diagnosis.government_pain_points:
        education_parts.append("Gov pain points: " + "; ".join(education_analysis.system_diagnosis.government_pain_points[:5]))

    strategy_parts = []
    if strategy.key_asks:
        strategy_parts.append("Key asks: " + "; ".join(strategy.key_asks[:5]))

    # Get LLM to synthesize state-specific content
    llm_content = await call_llm_plain(
        system_prompt=STATE_DECK_PROMPT.format(
            state=state,
            state_upper=state.upper(),
            state_data="\n".join(state_data_parts),
            education_data="\n".join(education_parts),
            strategy_summary="\n".join(strategy_parts),
        ),
        user_prompt=f"Produce the state-specific data for the {state} governor pitch deck.",
    )

    # Build the PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    _add_title_slide(prs, state)

    # --- Slide 2: The Problem ---
    problem_bullets = [
        f"{state} K-12 system: {k12_students:,.0f} students across the state",
        f"Per-pupil spending: ${per_pupil:,.0f}/student",
        "This isn't a funding problem — it's a model problem",
        "The traditional classroom model hasn't delivered results",
        "More money into the same system won't change outcomes",
    ]
    _add_content_slide(prs, state, f"The Problem — {state} Needs a New Model", problem_bullets)

    # --- Slide 3: We've Solved This ---
    results_bullets = [
        f"Alpha students: {ALPHA_RESULTS['avg_growth']} average growth (vs 1x expected)",
        f"Students 2 years behind: {ALPHA_RESULTS['behind_growth']} growth",
        f"{ALPHA_RESULTS['love_school']} of students love school (voted to keep school open over summer)",
        f"Avg SAT: {ALPHA_RESULTS['avg_sat']} vs 1063 national",
        f"Lipscomb Academy (Nashville): {ALPHA_RESULTS['lipscomb']}",
        f"Intervention pilots: ELA {ALPHA_RESULTS['intervention_ela']}, Math {ALPHA_RESULTS['intervention_math']}",
    ]
    _add_content_slide(prs, state, "We've Solved This — And We Can Prove It", results_bullets)

    # --- Slide 4: Live in Houston ---
    hisd_bullets = [
        "HISD pilot: America's largest public school pilot",
        "Partnership with Superintendent Miles",
        "2 campuses, 400 students, grades 3-8",
        "Alpha owns students >50% of the day",
        "Alpha investing $150k/campus + 5 guides per campus",
        "This proves it works in PUBLIC schools, not just private",
    ]
    _add_content_slide(prs, state, "Live in Houston — America's Largest Public School Pilot", hisd_bullets)

    # --- Slide 5: How It Works ---
    model_bullets = [
        "Children can master academics in 2 hours/day (Bloom's 2-Sigma)",
        "AI makes 1:1 tutoring possible at scale for the first time",
        "The freed time goes to life skills that parents actually want",
        "Traditional: 6 hrs classroom instruction",
        "Alpha Model: 2 hrs academic mastery + 4 hrs life skills",
        "Teachers become coaches/mentors — AI handles instruction",
    ]
    _add_content_slide(prs, state, "How It Works — The Reinvented School Day", model_bullets)

    # --- Slide 6: The Deal ---
    deal_bullets = [
        "Our commitment: We take the pilot risk. Outcomes-based. You pay nothing unless targets hit.",
        "Your commitment: When we hit numbers, you roll it out. Pre-agreed triggers.",
        "",
        f"INTERVENTION: ${PRICING['intervention']:,}/student/yr — 1-hr AI tutoring block",
        f"  → Trigger: 60%+ NWEA MAP Growth → expand to all low-performing schools",
        f"MARK ROBER SCIENCE: ${PRICING['mark_rober_science']:,}/student/yr — 1-hr science block",
        f"  → Trigger: hit targets → available to every public school",
        f"FULL TRANSFORMATION: ~${timeback_cost:,}/student/yr (20% of per-pupil)",
        f"  → Trigger: hit targets → aggressive expansion statewide",
    ]
    _add_content_slide(prs, state, "The Deal — We Prove It, You Scale It", deal_bullets)

    # --- Slide 7: The Economics ---
    econ_bullets = [
        f"{state} spends ${per_pupil:,.0f}/student. Here's what Alpha costs:",
        "",
        f"Intervention & Mark Rober Science: $2,000/student/year (outcomes-based)",
        f"  → {round(2000/per_pupil*100)}% of current per-pupil spend for transformational outcomes",
        f"Full School Transformation: ~${timeback_cost:,}/student (20% of per-pupil)",
        f"  → Remaining 80% (${per_pupil - timeback_cost:,.0f}) covers guides, life skills, facility, ops",
        f"  → Same total budget. World-class outcomes.",
    ]
    if esa_data and esa_data.get("esa_amount"):
        esa_amt = esa_data["esa_amount"]
        econ_bullets.append(f"")
        econ_bullets.append(f"ESA/Voucher: ${esa_amt:,}/student covers Alpha programs directly")
    _add_content_slide(prs, state, "The Economics — What It Costs", econ_bullets)

    # --- Slide 8: What Each Side Gets ---
    provides_bullets = [
        "WHAT ALPHA PROVIDES:",
        "  • Timeback AI platform",
        "  • AlphaCore life-skills curriculum",
        "  • Guide training (teacher coaching)",
        "  • Ongoing management + support",
        "  • Outcomes guarantee (intervention tier)",
        "",
        f"WHAT {state.upper()} PROVIDES:",
        "  • Access to schools/districts",
        "  • Existing per-pupil funding",
        "  • Teachers willing to adopt the model",
        "  • Political support for innovation",
        "  • Adherence to daily model",
    ]
    _add_content_slide(prs, state, f"What {state} Gets", provides_bullets)

    # --- Slide 9: Gated Expansion Roadmap ---
    roadmap_bullets = [
        "PHASE 1 — PROVE (SY 2026-27):",
        "  Pilot in low-performing schools | Alpha takes the risk | Outcomes-based",
        "  Trigger: 60%+ NWEA MAP Growth, third-party verified",
        "",
        "PHASE 2 — EXPAND (SY 2027-28):",
        "  All low-performing schools + willing districts | Existing per-pupil funding",
        "  Trigger: Results hold across pilot cohort",
        "",
        "PHASE 3 — TRANSFORM (SY 2028-29+):",
        "  Full school re-inventions statewide | Per-pupil + facility investment",
        "  Contract signed before pilot starts. Triggers pre-agreed.",
    ]
    _add_content_slide(prs, state, "The Roadmap — Gated Expansion", roadmap_bullets)

    # --- Slide 10: The Ask ---
    ask_bullets = [
        "We'll prove it works, at our financial risk",
        "You commit now that when we hit targets, it scales — fast",
        "",
        f"WHAT WE NEED FROM {state.upper()}:",
        "  1. Access to low-performing schools for fall 2026 pilot",
        "  2. A signed expansion commitment: hit triggers → statewide rollout",
        "  3. A champion in the governor's office to clear blockers",
        "",
        "WHAT YOU GET:",
        "  • Zero financial risk on the pilot",
        "  • Third-party verified results before any state money moves",
        f"  • {state} leads the nation in education innovation — that's the headline",
    ]
    _add_content_slide(prs, state, "The Ask", ask_bullets)

    # --- Final slide ---
    _add_closing_slide(prs, state)

    # Save
    output_dir = os.path.join(OUTPUT_DIR, state.lower().replace(" ", "_"))
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, f"{state.replace(' ', '_')}_governor_pitch.pptx")
    prs.save(path)
    logger.info("State deck saved: %s", path)
    return path


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, state: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK

    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"Alpha: AI-Powered Education for {state}"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Better outcomes. Lower cost. Proven at scale."
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "CONFIDENTIAL"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p3.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, state: str, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Content
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf2 = txBox.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets[:12]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = bullet
        # Emphasize section headers (ALL CAPS lines)
        if bullet and bullet == bullet.upper() and len(bullet) > 3 and not bullet.startswith("  "):
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = ACCENT
        elif bullet.startswith("  →") or bullet.startswith("  •"):
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        elif bullet.startswith("  "):
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(5)

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.text = f"CONFIDENTIAL — Alpha | {state}"
    fp.font.size = Pt(9)
    fp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _add_closing_slide(prs: Presentation, state: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK

    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI is transforming every industry."
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Education is last. The governor who moves first wins."
    p2.font.size = Pt(28)
    p2.font.color.rgb = ACCENT
    p2.bold = True
    p2.alignment = PP_ALIGN.CENTER
