"""Document Generation Agent — produces executive-quality investor deck and
investment memorandum (proposal document).

Post-workshop (March 16, 2026): Unified dual-school model for all sovereign nations.
No tiers. No PPP scaling. Operator & Licensor (Marriott model). 100/0 equity.
Output is a single combined document: investment memorandum with embedded deal terms.

Generates slide decks via Gamma API and DOCX with appendices, formatted for
C-suite / head-of-state audiences.  The DOCX proposal is structured as a full
investment memorandum modelled after Goldman Sachs / Morgan Stanley deal books.
Deck generation uses the Ed71 slide 8 format as the template reference
(see: Ed71_ The World's First AI-Native National Education System.pptx).
"""

from __future__ import annotations

import logging
import os
import asyncio
from datetime import datetime

from services.gamma import generate_and_wait, _extract_gamma_url, _extract_export_url

from docx import Document as DocxDocument
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn as docx_qn

from pptx import Presentation as PptxPresentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

from models.schemas import (
    CountryProfile, EducationAnalysis, Strategy,
    FinancialModel, FinancialAssumptions, AudienceType,
)
from services.llm import call_llm_plain
from config import OUTPUT_DIR

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Head-of-state preamble — prepended to every LLM prompt
# ---------------------------------------------------------------------------

HEAD_OF_STATE_PREAMBLE = """CRITICAL — EXTERNALLY PRESENTABLE ARTIFACTS:
These documents will be presented DIRECTLY to heads of state, sovereign rulers, royal family
members, and senior government ministers. The UAE version of these materials was personally
reviewed by the daughter of the #2 figure in Abu Dhabi and the #3 figure in the UAE.

ABSOLUTE RULES:
1. NEVER include internal classifications (tiers, scores, ratings, internal frameworks)
2. NEVER reference "Tier 1", "Tier 2", "Tier 3" or any scoring/ranking system
3. NEVER include analytical frameworks, scoring matrices, or internal assessment tools
4. NEVER reference PPP adjustments, GDP scaling, or pricing methodology
5. NEVER use phrases like "Alpha Relevance", "fit score", or internal jargon
6. Write in the voice of a trusted advisor presenting to royalty — measured, confident, specific
7. All financial figures are FIXED (the deal model does not change by country)
8. Research is for cultural context and narrative color only — it does NOT drive financial numbers
9. Use the country's formal diplomatic name and proper honorifics for leaders
10. Every artifact must be suitable for immediate external presentation without editing
"""

# ---------------------------------------------------------------------------
# Investment Memorandum Prompts (multi-section generation)
# ---------------------------------------------------------------------------

IM_SECTION_PROMPTS = {
    "executive_summary": HEAD_OF_STATE_PREAMBLE +
    """You are a Managing Director at Goldman Sachs' Investment Banking Division.
Write the EXECUTIVE SUMMARY section of an investment memorandum for a strategic education
partnership between Alpha Holdings, Inc. and {target}.

IMPORTANT DEAL MODEL:
- Operator & Licensor structure (Marriott hotel model) — NOT a joint venture
- Counterparty owns 100% of the local entity, Alpha owns 0% equity
- Alpha is the exclusive operator & licensor
- Dual-school model:
  Flagship: $40K-$100K tuition, capital city + biggest cities, 2-3 schools, 50% backstop
  National: FIXED $25K per-student budget, 100K student-year minimum commitment
- Fixed upfront development: $250M AlphaCore + $250M App R&D + $250M LifeSkills = $750M total (non-negotiable)
- Management fee: 10% of combined revenue (non-negotiable)
- Timeback license: 20% of combined revenue (non-negotiable)
- NOTE: The term sheet / deal terms should be embedded as a section within this document,
  producing ONE combined investment memorandum (not a separate deliverable).

Context:
{context}

Write 1,500-2,000 words covering:
- Transaction overview and key investment highlights (5-7 bullet points at the top)
- The opportunity: why this market, why now
- Brief description of the Alpha model
- Proposed Operator & Licensor structure (Marriott model) with 100/0 equity
- Dual-school model: Flagship schools + National schools
- Key financial highlights (revenue, EBITDA, IRR, MOIC)
- Risk/reward assessment summary
- Recommended next steps

Use formal investment banking prose. Include data points. Be specific and quantitative.
Research data provides color commentary and narrative context — financial figures are FIXED.""",
    "market_overview": HEAD_OF_STATE_PREAMBLE +
    """You are a VP at Goldman Sachs' Global Markets Research Division.
Write the MARKET OVERVIEW & MACRO ANALYSIS section of an investment memorandum for an
education partnership in {target}.

Context:
{context}

Country Data:
{country_data}

Write 2,000-2,500 words covering:
1. MACRO-ECONOMIC OVERVIEW — GDP, growth trajectory, per-capita income, wealth distribution,
   currency dynamics, sovereign ratings, fiscal position
2. DEMOGRAPHIC PROFILE — Population, youth bulge, urbanisation, middle-class growth,
   household income distribution
3. POLITICAL & REGULATORY LANDSCAPE — Government structure, stability, education governance,
   foreign investment regulations, public-private partnership frameworks
4. NATIONAL VISION & REFORM AGENDA — Vision plans, education modernisation targets,
   digital transformation initiatives
5. COMPETITIVE LANDSCAPE — Existing private education operators, international schools,
   edtech players, market gaps

Include tables where relevant (markdown format). Use formal analytical language.
Cite specific data points with numbers, years, and sources.""",
    "education_deep_dive": HEAD_OF_STATE_PREAMBLE +
    """You are a Senior Partner at McKinsey & Company's Education Practice.
Write the EDUCATION SECTOR DEEP DIVE section of an investment memorandum for {target}.

Context:
{context}

Education Data:
{education_data}

Write 2,000-2,500 words covering:
1. CURRENT STATE OF EDUCATION — System structure, enrollment, public vs private split,
   per-pupil spending, teacher workforce, outcomes (PISA, literacy, graduation rates)
2. PAIN POINTS & SYSTEM FAILURES — Student disengagement, outdated pedagogy,
   skills gap for AI economy, parent frustration, teacher burnout
3. REFORM LANDSCAPE — Active government reforms, budget allocations, timeline, key stakeholders
4. PRIVATE EDUCATION MARKET — Market size (TAM/SAM/SOM), growth rate, segmentation,
   pricing tiers, demand-supply gap
5. THE ALPHA OPPORTUNITY — Specific gaps that the Alpha model fills, competitive
   advantages vs incumbents, unique value propositions for this market

Include a market sizing table. Use consulting-quality analysis with specific data.""",
    "alpha_model": HEAD_OF_STATE_PREAMBLE +
    """You are the Chief Strategy Officer of Alpha Holdings, Inc.
Write the THE ALPHA MODEL section of an investment memorandum.

Context:
{context}

Write 1,500-2,000 words covering:
1. THE INNOVATION — How Timeback works (AI compresses core academics into 2 hours/day),
   what students do with the remaining time (STEM, sports, arts, entrepreneurship, life skills)
2. THE TECHNOLOGY STACK — AlphaCore (curriculum OS), Incept eduLLM (custom AI for local curriculum),
   Guide School (teacher transformation program), Learning Intelligence Platform
3. PROVEN OUTCOMES — UAE reference deal ($1.5B, 200K students), learning outcomes data,
   student engagement metrics, parent satisfaction
4. SCALABILITY — How the model adapts to different markets, cultures, curricula, and
   languages while maintaining core quality
5. COMPETITIVE MOAT — Proprietary AI, outcomes data flywheel, curriculum IP, teacher training IP,
   government relationships
6. LOCALIZATION FOR {target} — Specific adaptations required for this market
   (cultural, curriculum, language, regulatory)

Be specific about outcomes data and the UAE reference deal.""",
    "deal_structure": HEAD_OF_STATE_PREAMBLE +
    """You are a Partner at a top-tier M&A advisory firm.
Write the PROPOSED DEAL STRUCTURE & PARTNERSHIP MODEL section of an investment memorandum
for the Alpha Holdings, Inc. × {target} partnership.

IMPORTANT: This uses the Operator & Licensor model (Marriott hotel model) — NOT a JV.
- Counterparty owns 100% of the local entity, Alpha owns 0% equity
- Alpha is the exclusive operator & licensor
- Dual-school model:
  Flagship: $40K-$100K tuition, 2-3 schools in capital + biggest cities, 50% backstop
  National: FIXED $25K per-student budget, 100K student-year minimum commitment
- Fixed upfront development costs: $250M each (AlphaCore, App R&D, LifeSkills) = $750M total (non-negotiable)
- Management fee: 10% of combined revenue (non-negotiable)
- Timeback license: 20% of combined revenue (non-negotiable)
- Prepaid management + timeback fees scale by student count

Context:
{context}

Strategy Data:
{strategy_data}

Write 1,500-2,000 words covering:
1. PARTNERSHIP STRUCTURE — Operator & Licensor (Marriott model), rationale for 100/0 ownership,
   governance framework, board composition. NOT a JV — counterparty owns 100%, Alpha operates.
2. DUAL-SCHOOL MODEL — Flagship as proof-of-concept/brand anchor,
   National for scale. How the two school types interact strategically.
3. IP LICENSING & FEE STRUCTURE — FIXED $750M upfront development costs ($250M × 3),
   ongoing management fees (10%), Timeback license fees (20%), prepaid fee structure
4. CAPITAL STRUCTURE — Total capital requirement, deployment phasing, counterparty funds
   100% of entity, Alpha contributes operational expertise and IP
5. GOVERNANCE & CONTROLS — Decision rights, quality assurance, curriculum oversight,
   brand protection, exit provisions
6. RISK ALLOCATION — How risks are shared between parties, insurance, force majeure,
   regulatory change provisions

Include a fee structure summary table. Be specific about numbers.
NOTE: Financial figures are FIXED — do not derive them from country research data.""",
    "financial_analysis": HEAD_OF_STATE_PREAMBLE +
    """You are a Managing Director at Morgan Stanley's Financial Sponsors Group.
Write the FINANCIAL ANALYSIS section of an investment memorandum for the Alpha Holdings, Inc. × {target}
partnership.

IMPORTANT: Financial parameters are FIXED — do not derive from country GDP or PPP data.
- Flagship: $40K-$100K tuition (set by AGI of top 20% families)
- National: FIXED $25K per-student budget (non-negotiable)
- Fixed development costs: $250M each × 3 = $750M total (non-negotiable)
- Management fee: 10%, Timeback: 20% (both non-negotiable)
- Research data is "color commentary" for narrative context only

Context:
{context}

Financial Data:
{financial_data}

Write 2,500-3,000 words covering:
1. KEY ASSUMPTIONS — Enrollment ramp (Flagship + National), pricing
   (FIXED $25K national, $40K-$100K flagship), cost structure, capital expenditure,
   working capital. No PPP/GDP scaling — all financial figures are fixed.
2. 5-YEAR P&L PROJECTION — Year-by-year walkthrough of students, schools, revenue, COGS,
   gross margin, OPEX, EBITDA, net income, FCF. Include commentary on drivers and inflection points.
   Show Flagship and National revenue separately where relevant.
3. UNIT ECONOMICS — Per-student revenue, cost, and margin by school type
   (Flagship vs National). Compare to industry benchmarks.
4. REVENUE STREAMS — Management fees (10% of combined revenue), Timeback license (20%),
   FIXED $750M upfront development fees. Break down Alpha's operator revenue vs local entity revenue.
5. RETURNS ANALYSIS — IRR, MOIC, payback period, enterprise value at exit.
   Include sensitivity analysis on key variables.
6. CAPITAL DEPLOYMENT — Year-by-year capex schedule, FIXED $750M development costs,
   launch capital, real estate/infrastructure.
7. VALUATION & COMPARABLE TRANSACTIONS — How this deal compares to recent education sector
   M&A and partnerships globally. EV/EBITDA, EV/Revenue multiples.

Include P&L summary table, returns summary table, and sensitivity matrix. Use rigorous
financial analysis language.""",
    "implementation_roadmap": HEAD_OF_STATE_PREAMBLE +
    """You are a Senior Partner at Bain & Company's Private Equity Group.
Write the IMPLEMENTATION ROADMAP & EXECUTION PLAN section of an investment memorandum for
the Alpha Holdings, Inc. × {target} partnership.

Context:
{context}

Strategy Data:
{strategy_data}

Write 1,500-2,000 words covering:
1. PHASE 1: ESTABLISHMENT (Months 1-12) — Legal entity setup, regulatory approvals,
   site selection, initial hiring, curriculum localization, pilot school launch
2. PHASE 2: PROOF OF CONCEPT (Months 12-24) — First school operations, outcomes measurement,
   parent feedback, government engagement, second school pipeline
3. PHASE 3: SCALE-UP (Year 2-3) — Expansion to additional cities/regions, teacher training
   at scale, technology platform deployment, marketing ramp
4. PHASE 4: MARKET LEADERSHIP (Year 3-5) — Full portfolio deployment, brand establishment,
   government contract discussions, potential second-wave markets
5. KEY MILESTONES & DECISION GATES — Specific KPIs and go/no-go criteria for each phase
6. TEAM & ORGANIZATION — Key hires, organizational structure, advisory board, local leadership

Include a milestone timeline table.""",
    "risk_factors": HEAD_OF_STATE_PREAMBLE +
    """You are the Chief Risk Officer of a major investment bank.
Write the RISK FACTORS & MITIGATION section of an investment memorandum for the
Alpha Holdings, Inc. × {target} partnership.

Context:
{context}

Write 1,500-2,000 words covering:
1. REGULATORY & LICENSING RISK — Private school licensing, curriculum approval,
   foreign ownership restrictions, changing regulations. Mitigation strategies.
2. POLITICAL & SOVEREIGN RISK — Government change, policy reversals, geopolitical events.
   Mitigation through multi-stakeholder alignment.
3. EXECUTION & OPERATIONAL RISK — Teacher recruitment, technology deployment, construction,
   quality maintenance at scale. Mitigation through phased rollout.
4. FINANCIAL & FX RISK — Currency fluctuation, inflation, revenue collection, pricing pressure.
   Hedging strategies.
5. COMPETITIVE RISK — New entrants, incumbent responses, edtech disruption.
   Mitigation through proprietary moat.
6. CULTURAL & SOCIAL RISK — Community acceptance, cultural sensitivity, parental concerns.
   Mitigation through local advisory board and cultural IP layer.
7. TECHNOLOGY RISK — AI model reliability, data privacy, cybersecurity.
   Mitigation through robust tech governance.
8. REPUTATIONAL RISK — Student safety, outcomes delivery, media scrutiny.
   Mitigation through quality assurance framework.

Rate each risk (High/Medium/Low probability and impact). Include a risk matrix table.""",
    "appendices": HEAD_OF_STATE_PREAMBLE +
    """You are a VP at Goldman Sachs preparing the APPENDICES section of an
investment memorandum for the Alpha Holdings, Inc. × {target} partnership.

Context:
{context}

Financial Data:
{financial_data}

Write 1,500-2,000 words covering:
1. APPENDIX A: DETAILED FINANCIAL PROJECTIONS — Complete 5-year P&L table with all line items,
   detailed assumptions behind each number
2. APPENDIX B: COMPARABLE TRANSACTIONS — Table of 5-10 comparable education sector deals
   globally (date, buyer, target, size, multiple, structure)
3. APPENDIX C: REGULATORY CHECKLIST — Step-by-step regulatory requirements for establishing
   private schools in {target}
4. APPENDIX D: KEY TERMS SHEET — Summary of proposed commercial terms (fee schedule,
   ownership, governance, exit provisions)
5. APPENDIX E: REFERENCE CONTACTS — Key government contacts, potential local partners,
   advisory firms (use placeholder names)

Format as clean appendix sections with tables where appropriate.""",
}


async def generate_documents(
    target: str,
    country_profile: CountryProfile,
    education_analysis: EducationAnalysis,
    strategy: Strategy,
    financial_model: FinancialModel,
    assumptions: FinancialAssumptions,
    audience: AudienceType = AudienceType.INVESTOR,
    revision_notes: str | None = None,
    export_as: str = "pptx",
) -> tuple[str | None, str | None, str, str, str | None, str]:
    """Generate presentation deck (via Gamma with local fallback), investment memorandum, and spreadsheet.

    Returns (gamma_url, export_url, docx_path, xlsx_path, local_pptx_path, deck_input_text).
    ``local_pptx_path`` is a locally-generated PPTX fallback produced when
    the Gamma API is unavailable.  ``deck_input_text`` is the raw content
    sent to Gamma, enabling callers to request additional export formats.
    """
    logger.info("Generating documents for %s (audience: %s)", target, audience.value)

    output_dir = os.path.join(OUTPUT_DIR, target.lower().replace(" ", "_"))
    os.makedirs(output_dir, exist_ok=True)

    # --- Build context strings ---
    context = _build_context(country_profile, strategy, financial_model)
    country_data = _build_country_data(country_profile)
    education_data = _build_education_data(education_analysis)
    strategy_data = _build_strategy_data(strategy)
    financial_data = _build_financial_data(financial_model, assumptions)

    audience_labels = {
        AudienceType.ROYAL: "a Royal Family / Head of State",
        AudienceType.MINISTER: "a Minister of Education / Government Official",
        AudienceType.INVESTOR: "a VC/Institutional Investor",
    }

    # --- Get deck content from LLM ---
    deck_outline = await call_llm_plain(
        system_prompt=DECK_OUTLINE_PROMPT.format(
            target=target,
            context=context,
            audience_label=audience_labels.get(audience, "investor"),
            slide_count=11,
        ),
        user_prompt=(
            f"Produce the detailed slide outline for {target}."
            f"{' Additional notes: ' + revision_notes if revision_notes else ''}"
        ),
    )

    # --- Always generate local PPTX first (guaranteed fallback) ---
    local_pptx_path = None
    try:
        local_pptx_path = _build_pptx(
            target, strategy, financial_model, deck_outline, audience, output_dir,
        )
        logger.info("Local PPTX generated: %s", local_pptx_path)
    except Exception as exc:
        logger.error("Local PPTX generation failed for %s: %s", target, exc)

    # --- Generate investor deck via Gamma (enhanced version) ---
    gamma_url, export_url, deck_input_text = await _build_investor_deck_gamma(
        target, strategy, financial_model, deck_outline, audience,
        export_as=export_as,
    )

    # If Gamma succeeded, log it; otherwise the local PPTX is our backup
    if gamma_url or export_url:
        logger.info("Gamma deck generated successfully for %s", target)
    else:
        logger.warning(
            "Gamma unavailable for %s — local PPTX %s will be used as fallback",
            target, "is available" if local_pptx_path else "ALSO FAILED",
        )

    # --- Generate DOCX investment memorandum (multi-section) ---
    docx_path = await _build_investment_memorandum(
        target, country_profile, education_analysis, strategy,
        financial_model, assumptions, audience, output_dir,
        context, country_data, education_data, strategy_data, financial_data,
        revision_notes,
    )

    # --- Generate XLSX ---
    from agents.financial import export_model_xlsx
    xlsx_path = export_model_xlsx(target, financial_model, assumptions, country_profile)

    logger.info("Documents generated: gamma=%s, local_pptx=%s, docx=%s, xlsx=%s",
                gamma_url, local_pptx_path, docx_path, xlsx_path)
    return gamma_url, export_url, docx_path, xlsx_path, local_pptx_path, deck_input_text


# ---------------------------------------------------------------------------
# Context builders
# ---------------------------------------------------------------------------

def _build_context(cp: CountryProfile, strategy: Strategy, model: FinancialModel) -> str:
    parts = []
    if cp.demographics.total_population:
        parts.append(f"Population: {cp.demographics.total_population:,.0f}")
    if cp.demographics.population_0_18:
        parts.append(f"Population 0-18: {cp.demographics.population_0_18:,.0f}")
    if cp.economy.gdp:
        parts.append(f"GDP: ${cp.economy.gdp:,.0f}")
    if cp.economy.gdp_per_capita:
        parts.append(f"GDP/capita: ${cp.economy.gdp_per_capita:,.0f}")
    if cp.economy.gdp_growth_rate:
        parts.append(f"GDP growth: {cp.economy.gdp_growth_rate}%")
    if cp.economy.credit_rating:
        parts.append(f"Sovereign rating: {cp.economy.credit_rating}")
    if strategy.entry_mode:
        parts.append(f"Entry mode: {strategy.entry_mode.value}")
    if strategy.partnership_structure.type:
        parts.append(f"Partnership: {strategy.partnership_structure.type.value}")
    if strategy.partnership_structure.ownership_split:
        parts.append(f"Ownership: {strategy.partnership_structure.ownership_split}")
    if model.upfront_ip_fee:
        parts.append(f"Upfront IP fee: ${model.upfront_ip_fee:,.0f}")
    if model.management_fee_pct:
        parts.append(f"Management fee: {model.management_fee_pct * 100:.0f}%")
    if model.timeback_license_pct:
        parts.append(f"Timeback license: {model.timeback_license_pct * 100:.0f}%")
    if model.pnl_projection:
        for p in model.pnl_projection:
            parts.append(f"Year {p.year}: {p.students:,} students, {p.schools} schools, "
                         f"${p.revenue:,.0f} revenue, ${p.ebitda:,.0f} EBITDA, ${p.free_cash_flow:,.0f} FCF")
    if model.returns_analysis.irr:
        parts.append(f"IRR: {model.returns_analysis.irr}%")
    if model.returns_analysis.moic:
        parts.append(f"MOIC: {model.returns_analysis.moic}x")
    if model.returns_analysis.payback_period_years:
        parts.append(f"Payback: {model.returns_analysis.payback_period_years} years")
    if model.total_management_fee_revenue:
        parts.append(f"Total mgmt fee revenue (5yr): ${model.total_management_fee_revenue:,.0f}")
    if model.total_timeback_license_revenue:
        parts.append(f"Total Timeback license revenue (5yr): ${model.total_timeback_license_revenue:,.0f}")
    return "\n".join(parts)


def _build_country_data(cp: CountryProfile) -> str:
    parts = []
    d = cp.demographics
    e = cp.economy
    r = cp.regulatory
    pc = cp.political_context
    ed = cp.education

    parts.append(f"Target: {cp.target.name} ({cp.target.type.value})")
    if d.total_population:
        parts.append(f"Population: {d.total_population:,.0f}")
    if d.population_0_18:
        parts.append(f"Youth (0-18): {d.population_0_18:,.0f}")
    if d.growth_rate:
        parts.append(f"Pop growth: {d.growth_rate}%")
    if d.urbanisation:
        parts.append(f"Urbanisation: {d.urbanisation}%")
    if d.median_age:
        parts.append(f"Median age: {d.median_age}")
    if d.median_household_income:
        parts.append(f"Median HH income: ${d.median_household_income:,.0f}")
    if d.gini_coefficient:
        parts.append(f"Gini: {d.gini_coefficient}")
    if e.gdp:
        parts.append(f"GDP: ${e.gdp:,.0f}")
    if e.gdp_per_capita:
        parts.append(f"GDP/capita: ${e.gdp_per_capita:,.0f}")
    if e.gdp_growth_rate:
        parts.append(f"GDP growth: {e.gdp_growth_rate}%")
    if e.currency:
        parts.append(f"Currency: {e.currency}")
    if e.fx_rate:
        parts.append(f"FX rate: {e.fx_rate}")
    if e.inflation:
        parts.append(f"Inflation: {e.inflation}%")
    if e.sovereign_wealth_fund:
        parts.append(f"SWF: {e.sovereign_wealth_fund}")
    if e.credit_rating:
        parts.append(f"Rating: {e.credit_rating}")
    if r.ministry_of_education:
        parts.append(f"MoE: {r.ministry_of_education}")
    if r.foreign_ownership_rules:
        parts.append(f"Foreign ownership: {r.foreign_ownership_rules}")
    if r.ppp_framework:
        parts.append(f"PPP: {r.ppp_framework}")
    if pc.government_type:
        parts.append(f"Government: {pc.government_type}")
    if pc.head_of_state:
        parts.append(f"Head of state: {pc.head_of_state}")
    if pc.national_vision_plan:
        parts.append(f"Vision plan: {pc.national_vision_plan}")
    if ed.k12_enrolled:
        parts.append(f"K-12 enrolled: {ed.k12_enrolled:,.0f}")
    if ed.public_private_split:
        parts.append(f"Public/private: {ed.public_private_split}")
    if ed.avg_private_tuition:
        parts.append(f"Avg private tuition: ${ed.avg_private_tuition:,.0f}")
    return "\n".join(parts)


def _build_education_data(ea: EducationAnalysis) -> str:
    parts = []
    if ea.system_diagnosis.primary_pain_points:
        parts.append("Student pain points: " + "; ".join(ea.system_diagnosis.primary_pain_points[:5]))
    if ea.system_diagnosis.parent_pain_points:
        parts.append("Parent pain points: " + "; ".join(ea.system_diagnosis.parent_pain_points[:5]))
    if ea.system_diagnosis.government_pain_points:
        parts.append("Government pain points: " + "; ".join(ea.system_diagnosis.government_pain_points[:5]))
    if ea.reform_landscape.active_reforms:
        parts.append("Active reforms: " + "; ".join(ea.reform_landscape.active_reforms[:5]))
    if ea.reform_landscape.reform_budget:
        parts.append(f"Reform budget: {ea.reform_landscape.reform_budget}")
    if ea.reform_landscape.appetite_for_foreign_models:
        parts.append(f"Appetite for foreign models: {ea.reform_landscape.appetite_for_foreign_models}")
    if ea.two_hr_learning_fit.unique_value_propositions:
        parts.append("Alpha UVPs: " + "; ".join(ea.two_hr_learning_fit.unique_value_propositions[:5]))
    if ea.two_hr_learning_fit.model_recommendation:
        parts.append(f"Recommended entry: {ea.two_hr_learning_fit.model_recommendation.value}")
    return "\n".join(parts)


def _build_strategy_data(s: Strategy) -> str:
    parts = []
    if s.entry_mode:
        parts.append(f"Entry mode: {s.entry_mode.value}")
    if s.partnership_structure.type:
        parts.append(f"Partnership: {s.partnership_structure.type.value}")
    if s.partnership_structure.ownership_split:
        parts.append(f"Ownership: {s.partnership_structure.ownership_split}")
    if s.partnership_structure.ip_structure:
        parts.append(f"IP: {s.partnership_structure.ip_structure}")
    if s.brand.jv_name_suggestion:
        parts.append(f"Local entity name: {s.brand.jv_name_suggestion}")
    if s.school_types:
        for st in s.school_types[:4]:
            parts.append(f"School type: {st.name} — {st.focus or ''} — {st.tuition or ''}")
    if s.phased_rollout:
        for ph in s.phased_rollout[:5]:
            parts.append(f"Phase: {ph.phase} ({ph.timeline}) — {ph.student_count or 'TBD'} students")
    if s.key_asks:
        parts.append("Key asks: " + "; ".join(s.key_asks[:6]))
    if s.target_student_count_year5:
        parts.append(f"Y5 target: {s.target_student_count_year5:,} students")
    if s.upfront_ask:
        parts.append(f"Upfront ask: ${s.upfront_ask:,.0f}")
    return "\n".join(parts)


def _build_financial_data(model: FinancialModel, assumptions: FinancialAssumptions) -> str:
    parts = []
    for a in assumptions.assumptions:
        parts.append(f"{a.label}: {a.value} {a.unit}")
    parts.append("")
    if model.pnl_projection:
        parts.append("P&L Projection:")
        for p in model.pnl_projection:
            parts.append(
                f"  Y{p.year}: {p.students:,} students | {p.schools} schools | "
                f"${p.revenue:,.0f} rev | ${p.cogs:,.0f} COGS | ${p.gross_margin:,.0f} GM | "
                f"${p.opex:,.0f} OPEX | ${p.ebitda:,.0f} EBITDA | ${p.net_income:,.0f} NI | "
                f"${p.free_cash_flow:,.0f} FCF | ${p.cumulative_cash:,.0f} cum"
            )
    if model.unit_economics:
        parts.append("\nUnit Economics:")
        for ue in model.unit_economics:
            parts.append(f"  {ue.school_type}: ${ue.per_student_revenue:,.0f}/student, "
                         f"${ue.contribution_margin:,.0f} margin ({ue.margin_pct}%)")
    if model.returns_analysis:
        ra = model.returns_analysis
        parts.append(f"\nReturns: IRR {ra.irr}%, MOIC {ra.moic}x, Payback {ra.payback_period_years}y, "
                     f"EV@exit ${ra.enterprise_value_at_exit:,.0f}" if ra.enterprise_value_at_exit else "")
    if model.sensitivity:
        parts.append("\nSensitivity:")
        for s in model.sensitivity:
            parts.append(f"  {s.variable}: base={s.base_case}, down={s.downside}, up={s.upside}")
    if model.capital_deployment:
        parts.append("\nCapital Deployment:")
        for cd in model.capital_deployment:
            parts.append(f"  Y{cd.year}: ${cd.total:,.0f} total "
                         f"(IP ${cd.ip_development:,.0f}, Mgmt ${cd.management_fees:,.0f}, "
                         f"Launch ${cd.launch_capital:,.0f}, RE ${cd.real_estate:,.0f})")
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Investment Memorandum (DOCX) — multi-section generation
# ---------------------------------------------------------------------------

async def _build_investment_memorandum(
    target: str,
    country_profile: CountryProfile,
    education_analysis: EducationAnalysis,
    strategy: Strategy,
    financial_model: FinancialModel,
    assumptions: FinancialAssumptions,
    audience: AudienceType,
    output_dir: str,
    context: str,
    country_data: str,
    education_data: str,
    strategy_data: str,
    financial_data: str,
    revision_notes: str | None = None,
) -> str:
    """Build a comprehensive investment memorandum (30-50+ pages) with multiple LLM calls."""

    logger.info("Generating investment memorandum for %s...", target)

    sections: list[tuple[str, str]] = []  # (section_title, markdown_content)

    # Generate each section with a separate LLM call for maximum quality and length
    section_configs = [
        ("executive_summary", "I. EXECUTIVE SUMMARY"),
        ("market_overview", "II. MARKET OVERVIEW & MACRO ANALYSIS"),
        ("education_deep_dive", "III. EDUCATION SECTOR DEEP DIVE"),
        ("alpha_model", "IV. THE ALPHA MODEL"),
        ("deal_structure", "V. PROPOSED DEAL STRUCTURE & PARTNERSHIP MODEL"),
        ("financial_analysis", "VI. FINANCIAL ANALYSIS"),
        ("implementation_roadmap", "VII. IMPLEMENTATION ROADMAP & EXECUTION PLAN"),
        ("risk_factors", "VIII. RISK FACTORS & MITIGATION"),
        ("appendices", "IX. APPENDICES"),
    ]

    async def _generate_section(section_key: str, section_title: str) -> tuple[str, str]:
        prompt_template = IM_SECTION_PROMPTS[section_key]
        formatted_prompt = prompt_template.format(
            target=target,
            context=context,
            country_data=country_data,
            education_data=education_data,
            strategy_data=strategy_data,
            financial_data=financial_data,
        )
        user_prompt = f"Write the {section_title} section now."
        if revision_notes:
            user_prompt += f"\n\nAdditional revision notes from the user: {revision_notes}"
        try:
            content = await call_llm_plain(
                system_prompt=formatted_prompt,
                user_prompt=user_prompt,
            )
            logger.info("Generated section: %s (%d chars)", section_title, len(content))
            return (section_title, content)
        except Exception as exc:
            logger.error("Failed to generate section %s: %s", section_title, exc)
            return (section_title, f"[Section generation failed: {exc}]")

    sections = list(
        await asyncio.gather(
            *[_generate_section(key, title) for key, title in section_configs]
        )
    )

    # --- Build the DOCX ---
    doc = DocxDocument()

    # Page setup
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)

    # Set default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)
    style.font.color.rgb = DocxRGB(0x33, 0x33, 0x33)
    style.paragraph_format.space_after = DocxPt(6)
    style.paragraph_format.line_spacing = 1.15

    # Configure heading styles
    for level, (size, color, bold) in {
        1: (DocxPt(22), DocxRGB(0x1a, 0x1a, 0x2e), True),
        2: (DocxPt(16), DocxRGB(0x00, 0x6D, 0x77), True),
        3: (DocxPt(13), DocxRGB(0x33, 0x33, 0x33), True),
    }.items():
        heading_style = doc.styles[f"Heading {level}"]
        heading_style.font.size = size
        heading_style.font.color.rgb = color
        heading_style.font.bold = bold
        heading_style.font.name = "Calibri"
        heading_style.paragraph_format.space_before = DocxPt(18)
        heading_style.paragraph_format.space_after = DocxPt(8)

    # --- Cover Page ---
    doc.add_paragraph("")  # spacer
    doc.add_paragraph("")
    doc.add_paragraph("")

    cover_title = doc.add_paragraph()
    cover_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_title.add_run(f"Alpha Holdings, Inc. × {target}")
    run.font.size = DocxPt(36)
    run.font.color.rgb = DocxRGB(0x1a, 0x1a, 0x2e)
    run.bold = True

    cover_subtitle = doc.add_paragraph()
    cover_subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_subtitle.add_run("INVESTMENT MEMORANDUM")
    run.font.size = DocxPt(18)
    run.font.color.rgb = DocxRGB(0x00, 0x6D, 0x77)
    run.bold = True

    doc.add_paragraph("")

    cover_line = doc.add_paragraph()
    cover_line.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_line.add_run("Strategic Partnership Proposal for Education Transformation")
    run.font.size = DocxPt(14)
    run.font.color.rgb = DocxRGB(0x66, 0x66, 0x66)

    doc.add_paragraph("")
    doc.add_paragraph("")

    cover_conf = doc.add_paragraph()
    cover_conf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_conf.add_run("CONFIDENTIAL & PROPRIETARY")
    run.font.size = DocxPt(12)
    run.font.color.rgb = DocxRGB(0xCC, 0x00, 0x00)
    run.bold = True

    cover_date = doc.add_paragraph()
    cover_date.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_date.add_run(datetime.now().strftime("%B %Y"))
    run.font.size = DocxPt(12)
    run.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)

    doc.add_paragraph("")

    cover_prepared = doc.add_paragraph()
    cover_prepared.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover_prepared.add_run("Prepared by Alpha Holdings, Inc.")
    run.font.size = DocxPt(11)
    run.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)

    doc.add_page_break()

    # --- Disclaimer Page ---
    doc.add_heading("IMPORTANT NOTICE", level=2)
    disclaimer_text = (
        "This Investment Memorandum (the \"Memorandum\") has been prepared by Alpha Holdings, Inc. "
        "(\"Alpha\" or the \"Company\") solely for the purpose of providing information to "
        "prospective partners and investors in connection with a potential strategic partnership. "
        "This Memorandum is confidential and is being furnished to you solely for your information. "
        "It may not be reproduced or distributed to any other person.\n\n"
        "The information contained herein has been prepared in good faith and is believed to be "
        "reliable. However, no representation or warranty, express or implied, is made as to the "
        "accuracy, completeness, or fairness of the information and opinions contained in this "
        "Memorandum. Nothing in this document constitutes an offer, solicitation, or recommendation "
        "to invest.\n\n"
        "Certain statements in this Memorandum constitute forward-looking statements. These "
        "statements involve known and unknown risks, uncertainties, and other factors that may "
        "cause actual results to differ materially from those expressed or implied by such "
        "forward-looking statements.\n\n"
        "Recipients should conduct their own due diligence and seek independent professional "
        "advice before making any investment decision."
    )
    doc.add_paragraph(disclaimer_text)
    doc.add_page_break()

    # --- Table of Contents ---
    doc.add_heading("TABLE OF CONTENTS", level=1)
    doc.add_paragraph("")
    for _, section_title in section_configs:
        toc_entry = doc.add_paragraph()
        run = toc_entry.add_run(section_title)
        run.font.size = DocxPt(12)
        run.font.color.rgb = DocxRGB(0x00, 0x6D, 0x77)
        toc_entry.paragraph_format.space_after = DocxPt(6)

    doc.add_page_break()

    # --- Section Content ---
    for section_title, section_content in sections:
        doc.add_heading(section_title, level=1)

        # Parse markdown content into DOCX paragraphs
        _render_markdown_to_docx(doc, section_content)

        doc.add_page_break()

    # Save
    path = os.path.join(output_dir, f"{target.replace(' ', '_')}_investment_memorandum.docx")
    doc.save(path)
    logger.info("Investment memorandum saved: %s", path)
    return path


def _render_markdown_to_docx(doc: DocxDocument, md: str) -> None:
    """Parse markdown text and render into DOCX paragraphs, headings, lists, and tables."""
    lines = md.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            i += 1
            continue

        # Headings
        if stripped.startswith("#### "):
            doc.add_heading(stripped[5:], level=3)
            i += 1
            continue
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
            i += 1
            continue
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
            i += 1
            continue
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=2)  # Use level 2 since level 1 is section title
            i += 1
            continue

        # Tables (detect pipe-delimited markdown tables)
        if "|" in stripped and stripped.startswith("|"):
            table_lines = []
            while i < len(lines) and "|" in lines[i].strip() and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1

            # Filter out separator lines (e.g., |---|---|)
            data_lines = [line_text for line_text in table_lines if not all(c in "|-: " for c in line_text)]

            if len(data_lines) >= 1:
                # Parse cells
                rows = []
                for tl in data_lines:
                    cells = [c.strip() for c in tl.split("|") if c.strip()]
                    rows.append(cells)

                if rows:
                    ncols = max(len(r) for r in rows)
                    table = doc.add_table(rows=len(rows), cols=ncols)
                    table.style = "Light Grid Accent 1"
                    table.alignment = WD_TABLE_ALIGNMENT.CENTER

                    for row in table.rows:
                        tr = row._tr
                        tr_pr = tr.get_or_add_trPr()
                        tr_pr.append(tr_pr.makeelement(docx_qn("w:cantSplit"), {}))

                    for ri, row_data in enumerate(rows):
                        for ci, cell_text in enumerate(row_data):
                            if ci < ncols:
                                cell = table.rows[ri].cells[ci]
                                cell.text = cell_text
                                # Bold header row
                                if ri == 0:
                                    for paragraph in cell.paragraphs:
                                        for run in paragraph.runs:
                                            run.bold = True

                    doc.add_paragraph("")  # spacer after table
            continue

        # Bullet lists
        if stripped.startswith("- ") or stripped.startswith("* "):
            text = stripped[2:]
            # Check for bold prefix
            if text.startswith("**") and "**" in text[2:]:
                end_bold = text.index("**", 2)
                bold_part = text[2:end_bold]
                rest = text[end_bold + 2:]
                if rest.startswith(": ") or rest.startswith(" — ") or rest.startswith(" - "):
                    p = doc.add_paragraph(style="List Bullet")
                    r1 = p.add_run(bold_part)
                    r1.bold = True
                    p.add_run(rest)
                else:
                    doc.add_paragraph(text, style="List Bullet")
            else:
                doc.add_paragraph(text, style="List Bullet")
            i += 1
            continue

        # Numbered lists
        if len(stripped) >= 3 and stripped[0].isdigit() and (
                stripped[1] == "." or (stripped[1].isdigit() and stripped[2] == ".")):
            # Find the text after the number and period
            dot_idx = stripped.index(".")
            text = stripped[dot_idx + 1:].strip()
            doc.add_paragraph(text, style="List Number")
            i += 1
            continue

        # Horizontal rule
        if stripped in ("---", "***", "___"):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = DocxPt(12)
            p.paragraph_format.space_after = DocxPt(12)
            i += 1
            continue

        # Regular paragraph with bold/italic handling
        p = doc.add_paragraph()
        _add_formatted_text(p, stripped)
        i += 1


def _add_formatted_text(paragraph, text: str) -> None:
    """Add text to a paragraph with bold (**) and italic (*) formatting."""
    pos = 0
    while pos < len(text):
        # Bold
        bold_start = text.find("**", pos)
        if bold_start != -1:
            bold_end = text.find("**", bold_start + 2)
            if bold_end != -1:
                # Add text before bold
                if bold_start > pos:
                    paragraph.add_run(text[pos:bold_start])
                # Add bold text
                run = paragraph.add_run(text[bold_start + 2:bold_end])
                run.bold = True
                pos = bold_end + 2
                continue

        # No more formatting — add the rest
        paragraph.add_run(text[pos:])
        break


# ---------------------------------------------------------------------------
# PPTX Deck
# ---------------------------------------------------------------------------

DECK_OUTLINE_PROMPT = HEAD_OF_STATE_PREAMBLE + \
    """You are a senior director at Goldman Sachs' Investment Banking Division
drafting an investor/government presentation deck outline for a {audience_label} audience.

Target market: {target}
{context}

IMPORTANT: Use the Ed71 deck as the template reference for deal structure presentation.
Reference file: Ed71_ The World's First AI-Native National Education System.pptx

Deal model:
- Operator & Licensor — NOT a JV. Counterparty owns 100%, Alpha operates.
- Dual-school: Flagship ($40K-$100K) + National ($25K FIXED, 100K student-year min)
- Fixed development: $750M total ($250M × 3). Management fee 10%, Timeback 20%.

Produce a detailed slide-by-slide outline for an {slide_count}-slide presentation deck
covering:
1. Title slide with tagline
2. Vision — 'Educated in [Country]' as a Global Credential, Alpha as the Stanford of K-12 Education
3. The Key to Success — reinvented school day, core truths, Traditional vs Alpha Model comparison
4. The Alpha Model — Timeback, AlphaCore, Guide School, Incept eduLLM
5. Our Results — Bloom's 2-Sigma, growth metrics, SAT/AP scores, 97% love school, life skills projects
6. Market Opportunity — why this country, why now, key market drivers
7. The Complete Platform — Alpha education stack vs country-owned stack, deployment overview
8. School Type Portfolio — Halo Alpha School (flagship) and National school descriptions. National/country-owned schools MUST NOT have 'Alpha' in their name.
9. Country-Owned Schools & Investment — ownership structure, per-student budget, investment table
10. 5-Year Rollout Plan — phased approach with student targets
11. Thank You / closing

DO NOT include slides about: financial returns (IRR/MOIC), unit economics, risk mitigation,
key asks, P&L appendix, or capital deployment appendix.

For each slide provide: title, key bullet points (4-6), data callouts, and speaker notes.
"""


def _build_pptx(
    target: str,
    strategy: Strategy,
    model: FinancialModel,
    outline: str,
    audience: AudienceType,
    output_dir: str,
) -> str:
    """Build a local PPTX investor deck as a fallback when the Gamma API is
    unavailable.  Returns the absolute path to the generated ``.pptx`` file.

    The deck mirrors the same content structure used by ``_build_gamma_investor_input``
    but with professional tables, accent bars, and visual design elements.
    """
    from pptx.oxml.ns import qn

    prs = PptxPresentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Branding colors
    DARK_BG = PptxRGB(0x0A, 0x0F, 0x1A)
    ACCENT = PptxRGB(0x00, 0x6D, 0x77)
    ACCENT2 = PptxRGB(0x00, 0xD4, 0xAA)
    WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = PptxRGB(0xCC, 0xCC, 0xCC)
    MID_GRAY = PptxRGB(0x88, 0x88, 0x99)

    def _add_slide():
        layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        return slide

    def _add_accent_bar(slide, top: float = 0.0, height: float = 0.04):
        """Add an accent-colored bar across the top of the slide."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            PptxInches(0), PptxInches(top),
            PptxInches(13.333), PptxInches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()  # no border

    def _add_accent_line(slide, left: float, top: float, width: float):
        """Add a thin accent line for visual hierarchy."""
        shape = slide.shapes.add_shape(
            1,  # rectangle used as line
            PptxInches(left), PptxInches(top),
            PptxInches(width), PptxInches(0.03),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()

    def _add_slide_number(slide, num: int, total: int):
        """Add a slide number in the bottom-right."""
        txBox = slide.shapes.add_textbox(
            PptxInches(11.5), PptxInches(7.0), PptxInches(1.5), PptxInches(0.4),
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = f"{num} / {total}"
        p.font.size = PptxPt(9)
        p.font.color.rgb = MID_GRAY
        p.alignment = PP_ALIGN.RIGHT

    def _add_footer(slide):
        """Add confidential footer."""
        txBox = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(7.0), PptxInches(8.0), PptxInches(0.4),
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = "CONFIDENTIAL  |  Alpha Holdings, Inc."
        p.font.size = PptxPt(8)
        p.font.color.rgb = MID_GRAY
        p.alignment = PP_ALIGN.LEFT

    def _add_title(slide, text: str, top: float = 0.4):
        txBox = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(top), PptxInches(11.7), PptxInches(0.9),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    def _add_subtitle(slide, text: str, top: float = 1.2):
        txBox = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(top), PptxInches(11.7), PptxInches(0.5),
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = PptxPt(14)
        p.font.color.rgb = ACCENT
        p.font.italic = True

    def _add_body(slide, lines: list[str], top: float = 1.8):
        txBox = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(top), PptxInches(11.7), PptxInches(5.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = PptxPt(16)
            p.font.color.rgb = LIGHT_GRAY
            p.space_after = PptxPt(10)

    def _add_table(slide, headers: list[str], rows: list[list[str]],
                   left: float = 0.8, top: float = 1.8,
                   width: float = 11.7, row_height: float = 0.45):
        """Add a styled data table to the slide."""
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        tbl = slide.shapes.add_table(
            num_rows, num_cols,
            PptxInches(left), PptxInches(top),
            PptxInches(width), PptxInches(row_height * num_rows),
        ).table

        # Style header row
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = PptxPt(11)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            # Accent background
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn("a:solidFill"), {})
            srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "006D77"})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

        # Data rows
        for ri, row in enumerate(rows):
            bg_hex = "121624" if ri % 2 == 0 else "0E121E"
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = PptxPt(10)
                    p.font.color.rgb = LIGHT_GRAY
                    p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
                # Row background
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg_hex})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

        return tbl

    def _add_kpi_boxes(slide, kpis: list[tuple[str, str]], top: float = 1.8):
        """Add KPI highlight boxes in a row (label, value pairs)."""
        count = len(kpis)
        box_w = min(2.5, 11.0 / count)
        gap = 0.3
        total_w = count * box_w + (count - 1) * gap
        start_x = (13.333 - total_w) / 2

        for i, (label, value) in enumerate(kpis):
            x = start_x + i * (box_w + gap)
            # Box background
            shape = slide.shapes.add_shape(
                1, PptxInches(x), PptxInches(top),
                PptxInches(box_w), PptxInches(1.2),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptxRGB(0x12, 0x16, 0x24)
            shape.line.color.rgb = ACCENT
            shape.line.width = PptxPt(1)

            # Value
            txBox = slide.shapes.add_textbox(
                PptxInches(x + 0.1), PptxInches(top + 0.15),
                PptxInches(box_w - 0.2), PptxInches(0.6),
            )
            p = txBox.text_frame.paragraphs[0]
            p.text = value
            p.font.size = PptxPt(22)
            p.font.bold = True
            p.font.color.rgb = ACCENT2
            p.alignment = PP_ALIGN.CENTER

            # Label
            txBox2 = slide.shapes.add_textbox(
                PptxInches(x + 0.1), PptxInches(top + 0.75),
                PptxInches(box_w - 0.2), PptxInches(0.35),
            )
            p2 = txBox2.text_frame.paragraphs[0]
            p2.text = label
            p2.font.size = PptxPt(9)
            p2.font.color.rgb = MID_GRAY
            p2.alignment = PP_ALIGN.CENTER

    # Calculate total slides for numbering
    total_slides = 11
    slide_num = 0

    # ── Slide 1: Title ─────────────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s, top=0.0, height=0.06)
    _add_title(s, f"Alpha Holdings, Inc. × {target}", top=2.2)
    _add_body(s, ["Strategic Partnership Proposal"], top=3.3)
    # CONFIDENTIAL badge
    txBox = s.shapes.add_textbox(
        PptxInches(0.8), PptxInches(4.2), PptxInches(3.0), PptxInches(0.4),
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = "CONFIDENTIAL & NON-BINDING"
    p.font.size = PptxPt(11)
    p.font.bold = True
    p.font.color.rgb = PptxRGB(0xCC, 0x40, 0x40)
    # Date
    txBox2 = s.shapes.add_textbox(
        PptxInches(0.8), PptxInches(4.7), PptxInches(3.0), PptxInches(0.3),
    )
    from datetime import datetime as _dt
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = _dt.now().strftime("%B %Y")
    p2.font.size = PptxPt(10)
    p2.font.color.rgb = MID_GRAY
    # Bottom accent bar
    _add_accent_bar(s, top=7.44, height=0.06)

    # ── Slide 2: Executive Summary with KPI boxes ──────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Executive Summary")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    # KPI boxes
    kpis: list[tuple[str, str]] = []
    if model.pnl_projection:
        y5 = model.pnl_projection[-1]
        kpis.append(("Y5 Students", f"{y5.students:,}"))
        kpis.append(("Y5 Revenue", f"${y5.revenue:,.0f}"))
    if model.returns_analysis.irr:
        kpis.append(("IRR", f"{model.returns_analysis.irr}%"))
    if model.returns_analysis.moic:
        kpis.append(("MOIC", f"{model.returns_analysis.moic}x"))
    if kpis:
        _add_kpi_boxes(s, kpis[:4], top=1.5)

    exec_lines = [
        f"• Transform K-12 education in {target} through AI-powered learning",
        "• Operator & Licensor model (Marriott) — Counterparty owns 100%, Alpha operates",
        "• Dual-school: Flagship schools ($40K-$100K) + National schools ($25K fixed)",
        "• Proven model: UAE deal ($1.5B, 200K students) as reference",
    ]
    _add_body(s, exec_lines, top=3.2)

    # ── Slide 3: The Alpha Model ───────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "The Alpha Model")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_subtitle(s, "AI-powered education that compresses and elevates", top=1.3)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)
    _add_body(s, [
        "• Timeback: AI compresses core academics into 2 hours/day",
        "• Remaining time: STEM, sports, arts, entrepreneurship, life skills",
        "• AlphaCore: Curriculum OS managing the full student journey",
        "• Guide School: 12-month program transforming teachers into Guides",
        "• Incept eduLLM: Custom AI adapted to local curriculum & culture",
        "• Three commitments: Love school | Learn 2x faster | Future-ready skills",
    ])

    # ── Slide 4: Market Opportunity ────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, f"Market Opportunity: {target}")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)
    _add_body(s, [
        "• School-age population: significant K-12 cohort",
        "• Education sector undergoing reform and modernisation",
        "• Growing demand for premium, innovation-driven education",
        "• Gap between aspirations and current system performance",
        "• Government appetite for public-private partnerships",
        "• Alpha's model addresses the core pain points",
    ])

    # ── Slide 5: Financial Overview with TABLE ─────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "5-Year Financial Summary")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    if model.pnl_projection:
        headers = ["Metric"] + [f"Year {p.year}" for p in model.pnl_projection]
        pnl_rows = [
            ["Students"] + [f"{p.students:,}" for p in model.pnl_projection],
            ["Schools"] + [f"{p.schools}" for p in model.pnl_projection],
            ["Revenue"] + [f"${p.revenue:,.0f}" for p in model.pnl_projection],
            ["EBITDA"] + [f"${p.ebitda:,.0f}" for p in model.pnl_projection],
            ["Net Income"] + [f"${p.net_income:,.0f}" for p in model.pnl_projection],
            ["FCF"] + [f"${p.free_cash_flow:,.0f}" for p in model.pnl_projection],
        ]
        _add_table(s, headers, pnl_rows, top=1.6, row_height=0.42)
    else:
        _add_body(s, ["• Financial model pending"])

    # ── Slide 6: Returns & KPIs ────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Investment Returns")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    returns_kpis: list[tuple[str, str]] = [
        ("IRR", f"{model.returns_analysis.irr}%" if model.returns_analysis.irr else "N/A"),
        ("MOIC", f"{model.returns_analysis.moic}x" if model.returns_analysis.moic else "N/A"),
        ("Mgmt Fee (5yr)", f"${model.total_management_fee_revenue:,.0f}"),
        ("Timeback License (5yr)", f"${model.total_timeback_license_revenue:,.0f}"),
    ]
    _add_kpi_boxes(s, returns_kpis, top=1.5)
    _add_body(s, [
        f"• Fixed development costs: ${model.upfront_ip_fee:,.0f} ($250M × 3)",
        f"• Management fee: {model.management_fee_pct * 100:.0f}% of combined Flagship + National revenue",
        f"• Timeback license: {model.timeback_license_pct * 100:.0f}% of combined Flagship + National revenue",
    ], top=3.5)

    # ── Slide 7: Deal Structure ────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Proposed Deal Structure")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    deal_headers = ["Component", "Details"]
    deal_rows = [
        ["Structure", "Operator & Licensor (Marriott model)"],
        ["Ownership", "100/0 — Counterparty owns 100%, Alpha is exclusive operator & licensor"],
        ["Flagship", "$40K-$100K tuition, 2-3 schools, 50% backstop"],
        ["National", "$25K/student FIXED, 100K student-year min"],
        ["Development Costs", "$750M FIXED ($250M × 3)"],
        ["Management Fee", f"{model.management_fee_pct * 100:.0f}% of combined revenue"],
        ["Timeback License", f"{model.timeback_license_pct * 100:.0f}% of combined revenue"],
        ["Cultural IP Layer", "Local identity, language, and values fully integrated"],
    ]
    _add_table(s, deal_headers, deal_rows, top=1.6, row_height=0.5)

    # ── Slide 8: School Portfolio TABLE ────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "School Type Portfolio")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    if strategy.school_types and len(strategy.school_types) > 0:
        import re as _re
        school_headers = ["School Type", "Focus", "Tuition"]
        school_rows = []
        for st in strategy.school_types[:6]:
            name = st.name or ""
            is_flagship = any(kw in name.lower() for kw in ("flagship", "halo"))
            if not is_flagship and _re.search(r"\bAlpha\b", name):
                name = _re.sub(r"\s*\bAlpha\b\s*", " ", name).strip()
            school_rows.append([name, st.focus or "", st.tuition or ""])
        _add_table(s, school_headers, school_rows, top=1.6, row_height=0.5)
    else:
        _add_body(s, ["• Premium, Mid-Market, and Specialized school types"])

    # ── Slide 9: Rollout Plan ──────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "5-Year Rollout Plan")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    if strategy.phased_rollout:
        rollout_headers = ["Phase", "Timeline", "Students"]
        rollout_rows = [
            [ph.phase or "", ph.timeline or "", f"{ph.student_count:,}" if ph.student_count else ""]
            for ph in strategy.phased_rollout[:5]
        ]
        _add_table(s, rollout_headers, rollout_rows, top=1.6, row_height=0.5)
    else:
        _add_body(s, ["• Phased rollout details in strategy report"])

    # ── Slide 10: Unit Economics TABLE ─────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Unit Economics")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    if model.unit_economics:
        ue_headers = ["School Type", "Revenue/Student", "Margin/Student", "Margin %"]
        ue_rows = [
            [ue.school_type,
             f"${ue.per_student_revenue:,.0f}",
             f"${ue.contribution_margin:,.0f}",
             f"{ue.margin_pct}%"]
            for ue in model.unit_economics[:6]
        ]
        _add_table(s, ue_headers, ue_rows, top=1.6, row_height=0.5)
    else:
        _add_body(s, ["• Unit economics in financial model"])

    # ── Slide 11: Risk Mitigation ──────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Risk Mitigation")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)

    risk_headers = ["Risk Category", "Mitigation Strategy"]
    risk_rows = [
        ["Regulatory", "Proactive government engagement and compliance"],
        ["Execution", "Phased rollout with decision gates"],
        ["Cultural", "Local IP layer and cultural advisory board"],
        ["FX / Currency", "Local currency revenue with USD hedging strategy"],
        ["Competitive", "Proprietary AI and outcomes data as moat"],
        ["Political", "Multi-stakeholder alignment strategy"],
    ]
    _add_table(s, risk_headers, risk_rows, top=1.6, row_height=0.5)

    # ── Slide 12: Key Asks ─────────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s)
    _add_title(s, "Key Asks & Next Steps")
    _add_accent_line(s, 0.8, 1.15, 4.0)
    _add_slide_number(s, slide_num, total_slides)
    _add_footer(s)
    asks = strategy.key_asks[:6] if strategy.key_asks else [
        "Sovereign commitment to student volume targets",
        "Regulatory fast-track for school licensing",
        "Infrastructure/real estate support",
        "Cultural IP development partnership",
    ]
    _add_body(s, [f"• {a}" for a in asks])

    # ── Slide 13: Appendix – P&L (full table) ─────────────────────────
    if model.pnl_projection:
        slide_num += 1
        s = _add_slide()
        _add_accent_bar(s)
        _add_title(s, "Appendix: Detailed P&L Projection")
        _add_accent_line(s, 0.8, 1.15, 4.0)
        _add_slide_number(s, slide_num, total_slides)
        _add_footer(s)

        full_headers = ["Metric"] + [f"Year {p.year}" for p in model.pnl_projection]
        full_rows = [
            ["Students"] + [f"{p.students:,}" for p in model.pnl_projection],
            ["Schools"] + [f"{p.schools}" for p in model.pnl_projection],
            ["Revenue"] + [f"${p.revenue:,.0f}" for p in model.pnl_projection],
            ["COGS"] + [f"${p.cogs:,.0f}" for p in model.pnl_projection],
            ["Gross Margin"] + [f"${p.gross_margin:,.0f}" for p in model.pnl_projection],
            ["OPEX"] + [f"${p.opex:,.0f}" for p in model.pnl_projection],
            ["EBITDA"] + [f"${p.ebitda:,.0f}" for p in model.pnl_projection],
            ["Net Income"] + [f"${p.net_income:,.0f}" for p in model.pnl_projection],
            ["FCF"] + [f"${p.free_cash_flow:,.0f}" for p in model.pnl_projection],
        ]
        _add_table(s, full_headers, full_rows, top=1.6, row_height=0.4)

    # ── Slide 14: Appendix – Capital Deployment ───────────────────────
    if model.capital_deployment:
        slide_num += 1
        s = _add_slide()
        _add_accent_bar(s)
        _add_title(s, "Appendix: Capital Deployment")
        _add_accent_line(s, 0.8, 1.15, 4.0)
        _add_slide_number(s, slide_num, total_slides)
        _add_footer(s)

        cap_headers = ["Year", "IP Development", "Launch Capital", "Real Estate", "Total"]
        cap_rows = [
            [f"Year {cd.year}",
             f"${cd.ip_development:,.0f}" if cd.ip_development else "$0",
             f"${cd.launch_capital:,.0f}",
             f"${cd.real_estate:,.0f}",
             f"${cd.total:,.0f}"]
            for cd in model.capital_deployment
        ]
        _add_table(s, cap_headers, cap_rows, top=1.6, row_height=0.5)

    # ── Slide 15: Thank You ───────────────────────────────────────────
    slide_num += 1
    s = _add_slide()
    _add_accent_bar(s, top=0.0, height=0.06)
    _add_title(s, "Thank You", top=2.5)
    _add_body(s, [
        "Alpha Holdings, Inc. — Transforming Education Globally",
        "",
        "Contact: partnerships@2hrlearning.com",
    ], top=3.5)
    _add_accent_bar(s, top=7.44, height=0.06)

    # Save
    safe_name = target.lower().replace(" ", "_")
    filename = f"{safe_name}_investor_deck.pptx"
    pptx_path = os.path.join(output_dir, filename)
    prs.save(pptx_path)
    logger.info("Local PPTX saved: %s", pptx_path)
    return pptx_path


def _build_gamma_investor_input(
    target: str,
    strategy: Strategy,
    model: FinancialModel,
    outline: str,
    audience: AudienceType,
) -> str:
    """Build Gamma inputText for the investor deck with slide separators.

    11-slide structure aligned with the Ed71 reference deck.
    """
    slides: list[str] = []

    # --- Slide 1: Title ---
    slides.append(
        f"# Alpha Holdings, Inc. × {target}\n\n"
        f"Strategic Partnership Proposal\n\n"
        f"CONFIDENTIAL"
    )

    # --- Slide 2: Vision ---
    slides.append(
        f"# 'Educated in {target}' as a Global Credential\n\n"
        "Alpha — The \"Stanford of K-12 Education\"\n\n"
        "1. Only AI-native education system, purposefully designed for national scale\n"
        "2. Commitments to students: love school, learn 2× faster, life skills for the AI age\n"
        "3. Creating the next generation of global leaders through life skills and academic mastery\n\n"
        f"Our mission: Alpha is the AI-native public education system for all of {target}"
    )

    # --- Slide 3: The Key to Success ---
    slides.append(
        "# We have reinvented the school day\n\n"
        "Core truths of transformation:\n"
        "- Children should love school more than vacation\n"
        "- Children can master academics in 2 hours per day\n"
        "- The key to your children's happiness is high standards\n\n"
        "Traditional: 6 hours of classroom instruction\n"
        "Alpha Model: 2 hours academic mastery + 4 hours life-skills development\n\n"
        "Timeback: the AI and learning science platform delivering academic mastery "
        "10× faster than traditional schooling\n"
        "AlphaCore: an AI-age life-skills curriculum developing student leadership, "
        "teamwork, communication, resilience, and other non-academic capabilities"
    )

    # --- Slide 4: The Alpha Model ---
    slides.append(
        "# The Alpha Model\n\n"
        "Alpha reimagines the school day by compressing core academics into just two hours through AI, "
        "freeing the remaining time for STEM, sports, arts, entrepreneurship, and life skills. "
        "Three core commitments underpin every school: Love school | Learn 2× faster | Future-ready skills.\n\n"
        "Timeback: Proprietary AI engine that compresses core academic instruction into 2 hours per day "
        "without sacrificing depth or rigor, unlocking time for enrichment.\n"
        "AlphaCore: The Curriculum OS that manages the full student journey — from onboarding through assessment — "
        "ensuring consistency and continuous improvement across every school.\n"
        "Guide School: A 12-month transformation program that converts traditional teachers into Guides — "
        "facilitators of self-directed, AI-augmented learning experiences.\n"
        f"Incept eduLLM: A custom large language model adapted to the local curriculum, language, and cultural context "
        f"— ensuring the AI layer is authentically {target}, not merely translated."
    )

    # --- Slide 5: Our Results ---
    slides.append(
        "# The learning science has been known for 40 years… we have made it work\n\n"
        "Bloom's 2-Sigma Problem (1984): 1:1 tutoring produces 2 standard deviations of improvement. "
        "AI makes this possible at national scale for the first time.\n\n"
        "Exceptional Growth (vs. 1× expected): Avg all students 2.2×, Top 20% 3.9×, Top ⅓ 2.6×, 2 years behind 4.9×\n\n"
        "World-Class College Admissions: Avg. SAT 1530 (vs. 1063 national), 94% students with 4 or 5 on AP\n\n"
        "97% love school: High School students voted to keep school open over summer. "
        "Over 60% would rather go to school than on vacation. "
        "80% say their Guide is one of the most influential people in their life.\n\n"
        "100+ life skills projects: All Third Graders solve the Rubik's cube. "
        "Fifth Graders presented TED-style talks. Middle Schoolers placed 2nd in the world in Global AI Debates. "
        "High Schoolers traveled to Ukraine on a humanitarian mission. "
        "Students have created Broadway musicals, AI-powered mental health tools, and apps with millions of users."
    )

    # --- Slide 6: Market Opportunity ---
    slides.append(
        f"# Market Opportunity: {target}\n\n"
        f"Why {target}, Why Now:\n"
        f"- Significant K-12 school-age population providing a large addressable base\n"
        f"- Education sector actively undergoing reform and modernisation at the national level\n"
        f"- Growing demand for premium, innovation-driven educational alternatives\n"
        f"- Persistent gap between societal aspirations and current system performance\n"
        f"- Government appetite for structured public-private partnership models\n"
        f"- Alpha's model directly addresses the core pain points of the {target} system"
    )

    # --- Slide 7: The Complete Platform ---
    slides.append(
        f"# Launching {target}'s Alpha on the full education stack\n\n"
        "| Alpha Education Stack | "
        f"{target} Owned Stack |\n"
        "| --- | --- |\n"
        "| **Parent Education System:** Parents demand outcomes | "
        f"**Education Sovereignty:** {target} owns the critical pieces |\n"
        "| **Reinvented School Day:** 2 hrs academics + life skills | "
        "**Infrastructure:** Built to scale across 100+ schools |\n"
        "| **Timeback:** AI platform for 10× faster mastery | "
        f"**Localized AI Apps:** Built for {target} |\n"
        "| **AlphaCore:** Strongest K-12 life-skills curriculum | "
        "**Talent Academy:** Recruit and train local Guides |\n"
        "| **Guide School:** Talent academy for training Guides | "
        "**National eduLLM:** Embedded local laws, values, culture |\n"
        "| **Incept eduLLM:** Personalized content generation | "
        "**$25K/student** annual budget (~2× current public funding) |\n"
        "| **Scale:** 2,500 students → 100K across 100 campuses | "
        f"**First schools open SY26-27** in {target} |"
    )

    # --- Slide 8: School Type Portfolio ---
    if strategy.school_types:
        import re as _re
        school_lines = []
        for st in strategy.school_types[:4]:
            name = st.name or ""
            # National (non-flagship) school names must not contain "Alpha"
            is_flagship = any(kw in name.lower() for kw in ("flagship", "halo"))
            if not is_flagship and _re.search(r"\bAlpha\b", name):
                name = _re.sub(r"\s*\bAlpha\b\s*", " ", name).strip()
            school_lines.append(f"- {name}: {st.focus or ''} — {st.tuition or ''}")
        slides.append("# School Type Portfolio\n\n" + "\n".join(school_lines))
    else:
        slides.append(
            "# School Type Portfolio\n\n"
            "Halo Alpha School — $40K–$100K tuition:\n"
            "Innovative and personalized education. Flagship schools serve as the premium anchor — "
            "demonstrating the full Alpha experience at the highest level of execution.\n\n"
            "National School — $25K Fixed Budget:\n"
            "Accessible, high-quality education with personalized learning experiences at a fixed per-student cost. "
            "National schools are the engine of scale — targeting 100,000+ student-years."
        )

    # --- Slide 9: Country-Owned Schools & Investment ---
    slides.append(
        f"# {target} Schools: "
        f"{target} Owned, Alpha Operated\n\n"
        "*We are proposing to implement through a national network "
        "of privately-operated, government-funded schools, but are "
        "equally open to other structures.*\n\n"
        "**Partnership Structure:**\n"
        f"- 100% {target} owned, 0% Alpha owned. "
        "Alpha operates on behalf of the Country/State.\n"
        "- Per student funding/tuition: **$25,000/year**\n"
        "- Minimum **100,000 students per year** commitment\n"
        "- Schools operated as either public or private\n"
        f"- {target} sources real estate; schools pay rent\n\n"
        "**Investment Required:**\n\n"
        "| Category | Item | Amount |\n"
        "| --- | --- | --- |\n"
        "| Upfront Fixed | Alpha Core License | $250M |\n"
        "| Upfront Fixed | Country-Specific Incept EdLLM "
        "| $250M |\n"
        "| Upfront Fixed | Country-Specific Programs "
        "& Life Skills | $250M |\n"
        "| Upfront Fixed | Country-Specific EdTech Apps "
        "| $250M |\n"
        "| Prepaid | Timeback License Prepay | $500M |\n"
        "| Prepaid | Operating Fee Prepay | $250M |\n"
        "| **TOTAL UPFRONT** | | **$1.75B** |\n\n"
        "**Ongoing Annual:** Parent Education / Launch / "
        "Guides · Timeback License Fee (20% of funding) · "
        "Operating Fee (10% of funding)"
    )

    # --- Slide 10: 5-Year Rollout Plan ---
    if strategy.phased_rollout:
        phases = strategy.phased_rollout[:3]
        phase_rows = []
        for ph in phases:
            students = (
                f"{ph.student_count:,} students"
                if ph.student_count else "TBD"
            )
            phase_rows.append(
                f"| **{ph.phase}** | {ph.timeline} "
                f"| {students} |"
            )
        slides.append(
            "# 5-Year Rollout Plan\n\n"
            f"{target} scales systematically — "
            "building operational excellence before expanding, "
            "ensuring every cohort of students receives the "
            "full Alpha experience from day one.\n\n"
            "| Phase | Timeline | Target |\n"
            "| --- | --- | --- |\n"
            + "\n".join(phase_rows)
        )
    else:
        slides.append(
            "# 5-Year Rollout Plan\n\n"
            "Phased rollout details in strategy report"
        )

    # --- Slide 11: Thank You ---
    slides.append(
        "# Thank You\n\n"
        "Alpha Holdings, Inc. — Transforming Education Globally\n\n"
        "CONFIDENTIAL"
    )

    return "\n---\n".join(slides)


async def _build_investor_deck_gamma(
    target: str,
    strategy: Strategy,
    model: FinancialModel,
    outline: str,
    audience: AudienceType,
    export_as: str = "pptx",
) -> tuple[str | None, str | None, str]:
    """Build the investor deck via Gamma API.

    Returns (gamma_url, export_url, deck_input_text).
    The deck_input_text is returned so callers can request additional
    export formats (e.g. PDF) from Gamma if needed.

    Args:
        export_as: 'pptx' or 'pdf' — controls the format of the export URL.
    """
    input_text = _build_gamma_investor_input(target, strategy, model, outline, audience)

    try:
        result = await generate_and_wait(
            input_text,
            num_cards=11,
            text_mode="condense",
            card_split="inputTextBreaks",
            text_amount="extensive",
            additional_instructions=(
                f"This is a strategic partnership proposal / investor deck for {target}. "
                "The audience is C-suite / head-of-state level. "
                "Use a professional, data-driven tone. Keep slides clean with clear hierarchy. "
                "Use the markdown headings (# Title) as card titles. "
                "Preserve all financial figures, percentages, and data points exactly as provided. "
                "Do NOT generate charts, graphs, or bar charts. Use tables and text layouts only. "
                "On selected slides (not every slide), include high-quality images of modern schools, "
                "happy children learning, and innovative classroom environments. These should be aspirational "
                "and aligned with a premium education brand. Use icons where helpful for visual hierarchy. "
                "IMPORTANT: Country/state-owned schools (National) MUST NOT have 'Alpha' in their name. "
                "Only Halo Alpha Schools (Flagship, which are 100% Alpha-owned) may use the Alpha brand."
            ),
            export_as=export_as,
        )
    except Exception as exc:
        logger.error("Gamma API failed after retries for investor deck (%s): %s", target, exc)
        return None, None, input_text

    # Use robust URL extraction (handles multiple key name variations)
    gamma_url = _extract_gamma_url(result)
    export_url = _extract_export_url(result)

    if not gamma_url and not export_url:
        logger.error(
            "Gamma generation completed for %s but NO URLs found in response. "
            "Keys present: %s. Full response: %s",
            target, list(result.keys()), result,
        )
    else:
        logger.info("Investor deck generated via Gamma: url=%s, export=%s", gamma_url, export_url)

    return gamma_url, export_url, input_text
