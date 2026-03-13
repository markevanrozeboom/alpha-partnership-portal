"""PDF Generator — converts DOCX files to clean PDF documents.

Uses fpdf2 (pure Python, no system dependencies) to parse DOCX content
and render it as a professional PDF suitable for external customers.
"""

from __future__ import annotations

import logging
import os
import re
import unicodedata
from datetime import datetime

from docx import Document as DocxDocument
from fpdf import FPDF

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Colour palette (matching brand)
# ---------------------------------------------------------------------------
DARK = (26, 26, 46)          # #1a1a2e
ACCENT = (0, 109, 119)       # #006D77
WHITE = (255, 255, 255)
LIGHT_GRAY = (240, 240, 240)
MID_GRAY = (100, 100, 100)
TEXT_COLOR = (51, 51, 51)     # #333333
RED = (204, 0, 0)


def _clean_text(text: str) -> str:
    """Clean text for PDF rendering — handle encoding issues.

    The built-in Helvetica font in fpdf2 only supports latin-1 characters.
    This function aggressively normalises all Unicode so no unsupported
    character ever reaches the PDF renderer.
    """
    if not text:
        return text

    # Replace common problematic characters explicitly
    replacements: dict[str, str] = {
        "\u2018": "'",   # left single quote
        "\u2019": "'",   # right single quote
        "\u201a": "'",   # single low-9 quote
        "\u201c": '"',   # left double quote
        "\u201d": '"',   # right double quote
        "\u201e": '"',   # double low-9 quote
        "\u2013": "-",   # en dash
        "\u2014": "-",   # em dash
        "\u2015": "-",   # horizontal bar
        "\u2012": "-",   # figure dash
        "\u2026": "...", # ellipsis
        "\u00a0": " ",   # non-breaking space
        "\u2002": " ",   # en space
        "\u2003": " ",   # em space
        "\u2009": " ",   # thin space
        "\u200a": " ",   # hair space
        "\u200b": "",    # zero-width space
        "\u200c": "",    # zero-width non-joiner
        "\u200d": "",    # zero-width joiner
        "\ufeff": "",    # BOM / zero-width no-break space
        "\u2022": "-",   # bullet
        "\u2023": ">",   # triangular bullet
        "\u25aa": "-",   # black small square
        "\u25cb": "o",   # white circle
        "\u00b7": "-",   # middle dot
        "\u2010": "-",   # hyphen
        "\u2011": "-",   # non-breaking hyphen
        "\u00ab": '"',   # left guillemet
        "\u00bb": '"',   # right guillemet
        "\u2039": "'",   # single left angle quote
        "\u203a": "'",   # single right angle quote
        "\u2032": "'",   # prime
        "\u2033": '"',   # double prime
        "\u00ad": "",    # soft hyphen
        "\u00d7": "x",   # multiplication sign
        "\u2192": "->",  # rightwards arrow
        "\u2190": "<-",  # leftwards arrow
        "\u2194": "<->", # left right arrow
        "\u2713": "[x]", # check mark
        "\u2717": "[ ]", # ballot x
        "\u20ac": "EUR", # euro sign
        "\u00a3": "GBP", # pound sign
        "\u00a5": "JPY", # yen sign
    }
    for old, new in replacements.items():
        text = text.replace(old, new)

    # Normalise remaining Unicode via NFKD decomposition — this converts
    # accented chars into base + combining mark, then we keep only the
    # characters that survive latin-1 encoding.
    text = unicodedata.normalize("NFKD", text)

    # Final fallback: encode to latin-1, replacing anything still unsupported
    return text.encode("latin-1", errors="replace").decode("latin-1")


class AlphaPDF(FPDF):
    """Custom FPDF subclass with Alpha branding.

    All text output methods are wrapped so every string is automatically
    cleaned via ``_clean_text`` before reaching the PDF renderer.
    """

    def __init__(self, title: str = "", subtitle: str = ""):
        super().__init__()
        self.doc_title = _clean_text(title)
        self.doc_subtitle = _clean_text(subtitle)
        self.set_auto_page_break(auto=True, margin=20)

    # --- Wrap cell / multi_cell so callers never need to remember to clean ---

    def cell(self, w=None, h=None, text="", *args, **kwargs):
        return super().cell(w, h, _clean_text(str(text)), *args, **kwargs)

    def multi_cell(self, w=None, h=None, text="", *args, **kwargs):
        return super().multi_cell(w, h, _clean_text(str(text)), *args, **kwargs)

    # --- Branded header / footer ---

    def header(self):
        if self.page_no() == 1:
            return  # Skip header on cover page
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*MID_GRAY)
        super().cell(0, 8, _clean_text(
            f"CONFIDENTIAL - 2hr Learning (Alpha) | {self.doc_subtitle}"
        ), align="L")
        self.ln(4)
        # Thin accent line
        self.set_draw_color(*ACCENT)
        self.set_line_width(0.3)
        self.line(10, self.get_y(), self.w - 10, self.get_y())
        self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*MID_GRAY)
        super().cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

    def add_cover_page(self, title: str, subtitle: str, target: str):
        """Add a professional cover page."""
        self.add_page()
        # Dark background rectangle
        self.set_fill_color(*DARK)
        self.rect(0, 0, self.w, self.h, "F")

        # Main title
        self.set_y(80)
        self.set_font("Helvetica", "B", 36)
        self.set_text_color(*WHITE)
        self.multi_cell(0, 16, title, align="C")

        # Subtitle
        self.ln(8)
        self.set_font("Helvetica", "B", 18)
        self.set_text_color(*ACCENT)
        self.multi_cell(0, 10, subtitle, align="C")

        # Target
        self.ln(6)
        self.set_font("Helvetica", "", 14)
        self.set_text_color(170, 170, 170)
        self.multi_cell(0, 8, target, align="C")

        # Confidential notice
        self.ln(20)
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(204, 0, 0)
        self.multi_cell(0, 8, "CONFIDENTIAL & NON-BINDING", align="C")

        # Date
        self.ln(6)
        self.set_font("Helvetica", "", 11)
        self.set_text_color(150, 150, 150)
        self.multi_cell(0, 8, datetime.now().strftime("%B %Y"), align="C")

        # Prepared by
        self.ln(12)
        self.set_font("Helvetica", "", 10)
        self.set_text_color(150, 150, 150)
        self.multi_cell(0, 8, "Prepared by 2hr Learning - Alpha Division", align="C")


def convert_docx_to_pdf(docx_path: str) -> str:
    """Convert a DOCX file to a professional PDF.

    Reads the DOCX structure (paragraphs, headings, lists, tables)
    and renders them as a clean PDF document.

    Returns the path to the generated PDF file.
    """
    if not docx_path or not os.path.exists(docx_path):
        raise FileNotFoundError(f"DOCX file not found: {docx_path}")

    logger.info("Converting DOCX to PDF: %s", docx_path)

    doc = DocxDocument(docx_path)

    # Derive the PDF path
    pdf_path = docx_path.rsplit(".", 1)[0] + ".pdf"

    # Detect document type from filename
    basename = os.path.basename(docx_path).lower()
    if "term_sheet" in basename:
        doc_type = "term_sheet"
        cover_title = "INDICATIVE TERM SHEET"
        cover_subtitle = "Strategic Education Partnership"
    elif "investment_memorandum" in basename or "proposal" in basename:
        doc_type = "proposal"
        cover_title = "INVESTMENT MEMORANDUM"
        cover_subtitle = "Strategic Partnership Proposal for Education Transformation"
    else:
        doc_type = "document"
        cover_title = "DOCUMENT"
        cover_subtitle = ""

    # Extract target name from filename
    target_name = os.path.basename(docx_path).split("_")[0:3]
    target_name = " ".join(target_name).replace("_", " ").title()

    # Create PDF
    pdf = AlphaPDF(title=cover_title, subtitle=target_name)
    pdf.alias_nb_pages()

    # Cover page
    pdf.add_cover_page(
        title=cover_title,
        subtitle=cover_subtitle,
        target=f"2hr Learning (Alpha) x {target_name}",
    )

    # Content pages
    pdf.add_page()
    pdf.set_text_color(*TEXT_COLOR)

    for para in doc.paragraphs:
        text = _clean_text(para.text.strip())
        if not text:
            pdf.ln(3)
            continue

        style_name = para.style.name if para.style else ""

        # --- Headings ---
        if style_name.startswith("Heading 1") or style_name == "Heading 1":
            pdf.ln(6)
            pdf.set_font("Helvetica", "B", 18)
            pdf.set_text_color(*DARK)
            pdf.multi_cell(0, 9, text)
            # Accent underline
            pdf.set_draw_color(*ACCENT)
            pdf.set_line_width(0.5)
            pdf.line(10, pdf.get_y() + 1, 80, pdf.get_y() + 1)
            pdf.ln(5)
            pdf.set_text_color(*TEXT_COLOR)

        elif style_name.startswith("Heading 2"):
            pdf.ln(5)
            pdf.set_font("Helvetica", "B", 14)
            pdf.set_text_color(*ACCENT)
            pdf.multi_cell(0, 8, text)
            pdf.ln(2)
            pdf.set_text_color(*TEXT_COLOR)

        elif style_name.startswith("Heading 3"):
            pdf.ln(3)
            pdf.set_font("Helvetica", "B", 12)
            pdf.set_text_color(*TEXT_COLOR)
            pdf.multi_cell(0, 7, text)
            pdf.ln(2)

        # --- List items ---
        elif style_name.startswith("List Bullet"):
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(*TEXT_COLOR)
            x = pdf.get_x()
            pdf.cell(8, 6, "-")
            pdf.multi_cell(0, 6, text)
            pdf.ln(1)

        elif style_name.startswith("List Number"):
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(*TEXT_COLOR)
            pdf.cell(8, 6, " ")
            pdf.multi_cell(0, 6, text)
            pdf.ln(1)

        # --- Regular paragraphs ---
        else:
            # Check for bold runs
            has_bold = any(run.bold for run in para.runs if run.text.strip())
            has_italic = any(run.italic for run in para.runs if run.text.strip())

            if has_bold and len(para.runs) == 1:
                pdf.set_font("Helvetica", "B", 10)
            elif has_italic and len(para.runs) == 1:
                pdf.set_font("Helvetica", "I", 9)
                pdf.set_text_color(*MID_GRAY)
            else:
                pdf.set_font("Helvetica", "", 10)

            pdf.multi_cell(0, 6, text)
            pdf.ln(2)
            pdf.set_text_color(*TEXT_COLOR)

    # --- Render tables ---
    for table in doc.tables:
        pdf.ln(4)
        num_cols = len(table.columns)
        if num_cols == 0:
            continue

        col_width = (pdf.w - 20) / num_cols

        for ri, row in enumerate(table.rows):
            # Header row gets accent background
            if ri == 0:
                pdf.set_fill_color(*ACCENT)
                pdf.set_text_color(*WHITE)
                pdf.set_font("Helvetica", "B", 9)
            elif ri % 2 == 0:
                pdf.set_fill_color(*LIGHT_GRAY)
                pdf.set_text_color(*TEXT_COLOR)
                pdf.set_font("Helvetica", "", 9)
            else:
                pdf.set_fill_color(*WHITE)
                pdf.set_text_color(*TEXT_COLOR)
                pdf.set_font("Helvetica", "", 9)

            row_height = 7
            for ci, cell in enumerate(row.cells):
                cell_text = _clean_text(cell.text.strip())[:60]  # Truncate long cells
                pdf.cell(col_width, row_height, cell_text, border=1, fill=True)

            pdf.ln(row_height)

        pdf.ln(4)

    # Save
    pdf.output(pdf_path)
    logger.info("PDF generated: %s", pdf_path)
    return pdf_path


class LandscapeDeckPDF(FPDF):
    """Landscape PDF that mimics a slide deck — one slide per page.

    Each page has a dark background with white/accent text, matching the
    brand of the local PPTX fallback deck.
    """

    def __init__(self, target_name: str = ""):
        super().__init__(orientation="L", unit="mm", format="A4")
        # A4 landscape: 297 x 210 mm
        self.target_name = _clean_text(target_name)
        self.set_auto_page_break(auto=False)

    # --- Wrap cell / multi_cell for text cleaning ---

    def cell(self, w=None, h=None, text="", *args, **kwargs):
        return super().cell(w, h, _clean_text(str(text)), *args, **kwargs)

    def multi_cell(self, w=None, h=None, text="", *args, **kwargs):
        return super().multi_cell(w, h, _clean_text(str(text)), *args, **kwargs)

    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica", "I", 7)
        self.set_text_color(*MID_GRAY)
        super().cell(0, 8, (
            f"CONFIDENTIAL  |  2hr Learning (Alpha)  |  "
            f"{self.target_name}  |  Page {self.page_no()}"
        ), align="C")

    def _dark_bg(self):
        """Fill the current page with the dark brand background."""
        self.set_fill_color(10, 15, 26)  # #0A0F1A
        self.rect(0, 0, self.w, self.h, "F")

    def _accent_bar(self, y: float, width: float = 60):
        """Draw a horizontal accent bar."""
        self.set_draw_color(*ACCENT)
        self.set_line_width(0.8)
        self.line(15, y, 15 + width, y)

    def add_cover_slide(self, title: str, subtitle: str, target: str):
        """Dark cover slide matching the PPTX style."""
        self.add_page()
        self._dark_bg()

        # Accent bar at top
        self.set_fill_color(*ACCENT)
        self.rect(0, 0, self.w, 3, "F")

        # Main title
        self.set_y(55)
        self.set_font("Helvetica", "B", 36)
        self.set_text_color(*WHITE)
        self.multi_cell(0, 16, title, align="C")

        # Subtitle
        self.ln(6)
        self.set_font("Helvetica", "", 18)
        self.set_text_color(*ACCENT)
        self.multi_cell(0, 10, subtitle, align="C")

        # Target
        self.ln(4)
        self.set_font("Helvetica", "", 13)
        self.set_text_color(170, 170, 170)
        self.multi_cell(0, 8, target, align="C")

        # CONFIDENTIAL
        self.ln(14)
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(204, 60, 60)
        self.multi_cell(0, 8, "CONFIDENTIAL & NON-BINDING", align="C")

        # Date + prepared by
        self.ln(6)
        self.set_font("Helvetica", "", 9)
        self.set_text_color(130, 130, 130)
        self.multi_cell(0, 6, datetime.now().strftime("%B %Y"), align="C")
        self.multi_cell(0, 6, "Prepared by 2hr Learning - Alpha Division", align="C")

    def add_content_slide(self, title: str, bullets: list[str],
                          slide_num: int = 0, total_slides: int = 0):
        """Standard content slide with title + bullet points."""
        self.add_page()
        self._dark_bg()

        # Top accent bar
        self.set_fill_color(*ACCENT)
        self.rect(0, 0, self.w, 2.5, "F")

        # Slide number badge (top-right)
        if slide_num and total_slides:
            self.set_font("Helvetica", "", 8)
            self.set_text_color(100, 100, 100)
            self.set_xy(self.w - 35, 6)
            super().cell(25, 6, f"{slide_num} / {total_slides}", align="R")

        # Title
        self.set_xy(15, 12)
        self.set_font("Helvetica", "B", 24)
        self.set_text_color(*WHITE)
        self.multi_cell(self.w - 30, 11, title)

        # Accent underline
        title_bottom = self.get_y() + 2
        self._accent_bar(title_bottom, width=80)

        # Bullet content
        self.set_y(title_bottom + 6)
        self.set_x(18)

        for bullet in bullets:
            if not bullet.strip():
                self.ln(3)
                continue

            # Detect sub-heading style (lines starting with ** or all-caps short lines)
            is_subheading = (
                (bullet.startswith("**") and bullet.endswith("**")) or
                (len(bullet) < 50 and bullet == bullet.upper() and not bullet.startswith("-"))
            )
            clean = bullet.strip("*").strip("-").strip()

            if is_subheading:
                self.ln(3)
                self.set_font("Helvetica", "B", 13)
                self.set_text_color(*ACCENT)
                self.set_x(18)
                self.multi_cell(self.w - 36, 7, clean)
                self.ln(1)
            else:
                self.set_font("Helvetica", "", 11)
                self.set_text_color(200, 200, 210)
                self.set_x(18)
                # Bullet prefix
                super().cell(6, 7, _clean_text("-"))
                self.multi_cell(self.w - 42, 7, clean)
                self.ln(1)

            # Safety: don't overflow the page
            if self.get_y() > self.h - 20:
                break

    def add_table_slide(self, title: str, headers: list[str],
                        rows: list[list[str]], slide_num: int = 0,
                        total_slides: int = 0):
        """Slide with a styled data table."""
        self.add_page()
        self._dark_bg()

        # Top accent bar
        self.set_fill_color(*ACCENT)
        self.rect(0, 0, self.w, 2.5, "F")

        # Slide number badge
        if slide_num and total_slides:
            self.set_font("Helvetica", "", 8)
            self.set_text_color(100, 100, 100)
            self.set_xy(self.w - 35, 6)
            super().cell(25, 6, f"{slide_num} / {total_slides}", align="R")

        # Title
        self.set_xy(15, 12)
        self.set_font("Helvetica", "B", 24)
        self.set_text_color(*WHITE)
        self.multi_cell(self.w - 30, 11, title)
        self._accent_bar(self.get_y() + 2, width=80)
        self.ln(8)

        # Table
        num_cols = len(headers)
        if num_cols == 0:
            return
        usable = self.w - 30
        col_w = usable / num_cols
        row_h = 9

        # Header row
        self.set_x(15)
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(*ACCENT)
        self.set_text_color(*WHITE)
        for h in headers:
            super().cell(col_w, row_h, _clean_text(h), border=0, fill=True, align="C")
        self.ln(row_h)

        # Data rows
        for ri, row in enumerate(rows):
            if self.get_y() > self.h - 22:
                break  # don't overflow
            self.set_x(15)
            if ri % 2 == 0:
                self.set_fill_color(18, 22, 36)
            else:
                self.set_fill_color(14, 18, 30)
            self.set_font("Helvetica", "", 9)
            self.set_text_color(200, 200, 210)
            for ci, val in enumerate(row):
                align = "L" if ci == 0 else "R"
                super().cell(col_w, row_h, _clean_text(val), border=0, fill=True, align=align)
            self.ln(row_h)


def convert_pptx_to_pdf(pptx_path: str) -> str:
    """Convert a PPTX presentation to a professional landscape PDF deck.

    Each slide becomes a full landscape page with the dark-branded style,
    matching the look of the local PPTX fallback.

    Returns the path to the generated PDF file.
    """
    if not pptx_path or not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    logger.info("Converting PPTX to landscape PDF deck: %s", pptx_path)

    from pptx import Presentation

    prs = Presentation(pptx_path)
    pdf_path = pptx_path.rsplit(".", 1)[0] + ".pdf"

    # Extract target name
    parts = os.path.basename(pptx_path).split("_")[0:2]
    target_name = " ".join(parts).replace("_", " ").title()

    total_slides = len(prs.slides)
    pdf = LandscapeDeckPDF(target_name=target_name)

    # Cover slide
    pdf.add_cover_slide(
        title=f"2hr Learning x {target_name}",
        subtitle="Strategic Partnership Proposal",
        target=f"Prepared for {target_name}",
    )

    # Content slides
    for slide_idx, slide in enumerate(prs.slides, 1):
        # Gather all text content from this slide
        title_text = ""
        body_lines: list[str] = []
        table_data: list[list[str]] = []
        table_headers: list[str] = []

        for shape in slide.shapes:
            # Check for tables
            if shape.has_table:
                tbl = shape.table
                for ri, row in enumerate(tbl.rows):
                    cells = [_clean_text(c.text.strip()) for c in row.cells]
                    if ri == 0:
                        table_headers = cells
                    else:
                        table_data.append(cells)
                continue

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect title-like text (large font or first textbox)
                is_title = False
                for run in para.runs:
                    if run.font.size and run.font.size.pt >= 20:
                        is_title = True
                        break

                if is_title and not title_text:
                    title_text = text
                else:
                    body_lines.append(text)

        # If no title found, use first body line
        if not title_text and body_lines:
            title_text = body_lines.pop(0)
        if not title_text:
            title_text = f"Slide {slide_idx}"

        # Render as table slide or content slide
        if table_headers and table_data:
            pdf.add_table_slide(
                title_text, table_headers, table_data,
                slide_num=slide_idx, total_slides=total_slides,
            )
            # Also render any remaining body lines
            if body_lines:
                y = pdf.get_y() + 4
                pdf.set_y(y)
                pdf.set_x(18)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(180, 180, 190)
                for line in body_lines[:4]:
                    pdf.set_x(18)
                    pdf.multi_cell(pdf.w - 36, 6, _clean_text(line))
                    pdf.ln(1)
        else:
            pdf.add_content_slide(
                title_text, body_lines,
                slide_num=slide_idx, total_slides=total_slides,
            )

    pdf.output(pdf_path)
    logger.info("PPTX -> landscape PDF deck generated: %s", pdf_path)
    return pdf_path
